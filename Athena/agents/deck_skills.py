"""
Latitude MedTech — Deck Skills
================================
python-pptx slide-type renderers for DeckAgent.
Theme: "Advancing Human Intelligence" — Latitude MedTech brand.

Canvas: 13.333 × 7.5 inches (16:9), all measurements in EMU.
Quality standard: McKinsey/PwC management consulting register.
"""
from __future__ import annotations
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Brand palette ─────────────────────────────────────────────────────────────
_NAVY  = RGBColor(0x0A, 0x25, 0x40)   # headers, cover, primary dark
_SLATE = RGBColor(0x2C, 0x3E, 0x50)   # body text
_BLUE  = RGBColor(0x1A, 0x6F, 0xA3)   # accents, number chips
_TEAL  = RGBColor(0x1F, 0x7A, 0x6D)   # section dividers, ISO accent
_WARM  = RGBColor(0xC4, 0x92, 0x2A)   # recommendation, call-outs
_MUTED = RGBColor(0x8A, 0x9B, 0xB0)   # captions, footer text
_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
_LIGHT = RGBColor(0xF4, 0xF5, 0xF7)   # alternating rows
_RULE  = RGBColor(0xE3, 0xE6, 0xEA)   # divider lines
_RED   = RGBColor(0xC0, 0x39, 0x2B)   # mistake / warning accent

FONT = "Calibri"

# Canvas
W      = Inches(13.333)
H      = Inches(7.5)
LM     = Inches(0.65)
RM     = Inches(0.65)
CW     = W - LM - RM
HDR_H  = Inches(0.95)      # content-slide header band
FOOT_Y = Inches(7.05)

_TAGLINE = "Advancing Human Intelligence"


class DeckSkills:
    """Append brand-styled slides to a python-pptx Presentation."""

    def __init__(self, prs: Presentation):
        self.prs = prs
        self._blank = self._find_blank_layout()
        self._slide_count = 0   # incremented per slide; used for page numbers

    # ── Internal helpers ──────────────────────────────────────────────────────

    def _find_blank_layout(self):
        for layout in self.prs.slide_layouts:
            if "blank" in layout.name.lower():
                return layout
        return self.prs.slide_layouts[-1]

    def _slide(self):
        self._slide_count += 1
        return self.prs.slides.add_slide(self._blank)

    def _rect(self, slide, left, top, width, height, color: RGBColor):
        shape = slide.shapes.add_shape(1, int(left), int(top), int(width), int(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        return shape

    def _tx(self, slide, text: str, left, top, width, height, *,
            bold=False, italic=False, size=15, color=_SLATE,
            align=PP_ALIGN.LEFT, wrap=True):
        box = slide.shapes.add_textbox(int(left), int(top), int(width), int(height))
        tf = box.text_frame
        tf.word_wrap = wrap
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.bold   = bold
        run.font.italic = italic
        run.font.size   = Pt(size)
        run.font.color.rgb = color
        run.font.name   = FONT
        return box

    def _header(self, slide, title: str, subtitle: str = ""):
        """Dark navy header band — white title — optional muted subtitle."""
        # Navy band
        self._rect(slide, 0, 0, W, HDR_H, _NAVY)
        # Thin teal rule at bottom of band
        self._rect(slide, 0, HDR_H - Emu(9000), W, Emu(9000), _TEAL)
        # Title (white, inside band)
        title_top = Inches(0.13) if not subtitle else Inches(0.08)
        title_h   = Inches(0.6) if not subtitle else Inches(0.48)
        self._tx(slide, title,
                 LM, title_top, CW, title_h,
                 bold=True, size=22, color=_WHITE)
        if subtitle:
            self._tx(slide, subtitle,
                     LM, Inches(0.60), CW, Inches(0.30),
                     size=11, color=RGBColor(0xA8, 0xC4, 0xD8))

    def _footer(self, slide, label: str = "CONFIDENTIAL — Steven Tran Review Required"):
        n = self._slide_count
        self._rect(slide, 0, FOOT_Y, W, Emu(6000), _RULE)
        self._tx(slide, f"Latitude MedTech LLC  ·  {label}",
                 LM, FOOT_Y + Emu(10000), Inches(9), Inches(0.3),
                 size=8, color=_MUTED)
        if n > 1:   # no page number on cover
            self._tx(slide, str(n - 1),   # slide 1 = cover (unnumbered); content starts at 1
                     W - Inches(1.2), FOOT_Y + Emu(10000), Inches(0.8), Inches(0.3),
                     size=8, color=_MUTED, align=PP_ALIGN.RIGHT)

    # ── Slide types ───────────────────────────────────────────────────────────

    def add_cover(self, title: str, subtitle: str = "",
                  date: str = "", client_name: str = "") -> object:
        """Dark full-bleed cover — Advancing Human Intelligence brand."""
        slide = self._slide()

        # Dark background + left accent strip
        self._rect(slide, 0, 0, W, H, _NAVY)
        self._rect(slide, 0, 0, Emu(42000), H, _BLUE)

        # Firm name — top right
        self._tx(slide, "LATITUDE MEDTECH LLC",
                 W - Inches(4.2), Inches(0.3), Inches(3.7), Inches(0.4),
                 size=10, bold=True, color=RGBColor(0x1A, 0x6F, 0xA3),
                 align=PP_ALIGN.RIGHT)

        # Tagline — just below firm name
        self._tx(slide, _TAGLINE,
                 W - Inches(4.2), Inches(0.62), Inches(3.7), Inches(0.3),
                 size=9, italic=True, color=_MUTED, align=PP_ALIGN.RIGHT)

        # Title
        self._tx(slide, title,
                 Inches(0.85), Inches(1.9), Inches(10.5), Inches(2.4),
                 bold=True, size=40, color=_WHITE)

        # Horizontal rule below title block
        self._rect(slide, Inches(0.85), Inches(4.45), Inches(6.5), Emu(12000), _BLUE)

        # Subtitle
        if subtitle:
            self._tx(slide, subtitle,
                     Inches(0.85), Inches(4.65), Inches(10.5), Inches(0.65),
                     size=20, color=RGBColor(0xA8, 0xC4, 0xD8))

        # Meta line (date · client)
        meta = "  ·  ".join(p for p in [date, client_name] if p)
        if meta:
            meta_y = Inches(5.4) if subtitle else Inches(4.75)
            self._tx(slide, meta,
                     Inches(0.85), meta_y, Inches(9), Inches(0.35),
                     size=12, color=_MUTED)

        # Bottom confidential
        self._tx(slide, "CONFIDENTIAL",
                 W - Inches(2.8), Inches(7.05), Inches(2.3), Inches(0.35),
                 size=9, color=_MUTED, align=PP_ALIGN.RIGHT)
        return slide

    def add_exec_summary(self, situation: str, complication: str,
                         recommendation: str, findings: list[str]) -> object:
        """Executive summary — SCR left column, key findings right."""
        slide = self._slide()
        self._header(slide, "Executive Summary")
        self._footer(slide)

        col_w   = Inches(5.9)
        sep_x   = LM + col_w + Inches(0.3)
        right_x = sep_x + Inches(0.25)
        right_w = W - right_x - RM
        top     = HDR_H + Inches(0.18)

        scr = [("SITUATION",      situation,      _BLUE),
               ("COMPLICATION",   complication,   _WARM),
               ("RECOMMENDATION", recommendation, _TEAL)]
        y = top
        for label, body, accent in scr:
            self._rect(slide, LM, y, Emu(15000), Inches(0.72), accent)
            self._tx(slide, label,
                     LM + Inches(0.2), y + Emu(10000),
                     col_w - Inches(0.3), Inches(0.28),
                     bold=True, size=8, color=accent)
            self._tx(slide, body,
                     LM + Inches(0.2), y + Emu(28000),
                     col_w - Inches(0.3), Inches(0.85),
                     size=13, color=_SLATE, wrap=True)
            y += Inches(1.22)

        # Vertical rule
        self._rect(slide, sep_x, top, Emu(7000), Inches(3.7), _RULE)

        # Findings column
        self._tx(slide, "KEY FINDINGS",
                 right_x, top, right_w, Inches(0.28),
                 bold=True, size=9, color=_NAVY)
        y_r = top + Inches(0.38)
        for i, f in enumerate(findings[:5]):
            chip_x = right_x
            self._rect(slide, chip_x, y_r + Emu(14000),
                       Emu(48000), Emu(48000), _BLUE)
            self._tx(slide, f"{i + 1:02d}",
                     chip_x + Emu(8000), y_r + Emu(10000),
                     Emu(48000), Emu(52000),
                     bold=True, size=11, color=_WHITE, align=PP_ALIGN.CENTER)
            self._tx(slide, f,
                     right_x + Inches(0.58), y_r,
                     right_w - Inches(0.65), Inches(0.75),
                     size=13, color=_SLATE, wrap=True)
            y_r += Inches(0.85)
        return slide

    def add_section_divider(self, section_title: str, subtitle: str = "") -> object:
        """Full-bleed dark section break."""
        slide = self._slide()
        self._rect(slide, 0, 0, W, H, _NAVY)
        self._rect(slide, 0, 0, Emu(42000), H, _TEAL)

        self._tx(slide, "SECTION",
                 Inches(0.85), Inches(0.5), Inches(4), Inches(0.3),
                 size=9, color=_MUTED)
        self._tx(slide, section_title,
                 Inches(0.85), Inches(2.3), Inches(11.5), Inches(1.5),
                 bold=True, size=42, color=_WHITE)
        if subtitle:
            self._tx(slide, subtitle,
                     Inches(0.85), Inches(3.95), Inches(11), Inches(0.55),
                     size=19, color=RGBColor(0xA8, 0xC4, 0xD8))
        self._rect(slide, Inches(0.85), Inches(4.9), Inches(3.5), Emu(10000), _TEAL)

        # Tagline bottom right
        self._tx(slide, _TAGLINE,
                 W - Inches(4.5), Inches(6.95), Inches(4.2), Inches(0.35),
                 size=9, italic=True, color=_MUTED, align=PP_ALIGN.RIGHT)
        return slide

    def add_content_bullets(self, headline: str, bullets: list[str],
                            note: str = "") -> object:
        """Standard consulting body slide — headline + up to 5 bullets."""
        slide = self._slide()
        self._header(slide, headline)
        self._footer(slide)

        y = HDR_H + Inches(0.22)
        for bullet in bullets[:5]:
            # Thin left accent strip
            self._rect(slide, LM, y + Emu(28000), Emu(10000), Inches(0.52), _BLUE)
            self._tx(slide, bullet,
                     LM + Inches(0.18), y,
                     CW - Inches(0.22), Inches(0.82),
                     size=15, color=_SLATE, wrap=True)
            y += Inches(1.06)

        if note:
            self._tx(slide, f"Note: {note}",
                     LM, Inches(6.65), CW, Inches(0.28),
                     size=8, italic=True, color=_MUTED)
        return slide

    def add_two_column(self, headline: str,
                       left_title: str, left_bullets: list[str],
                       right_title: str, right_bullets: list[str]) -> object:
        """Side-by-side comparison slide."""
        slide = self._slide()
        self._header(slide, headline)
        self._footer(slide)

        half_w  = (CW - Inches(0.3)) / 2
        right_x = LM + half_w + Inches(0.3)
        top     = HDR_H + Inches(0.2)

        self._rect(slide, LM,      top, half_w, Inches(0.44), _BLUE)
        self._rect(slide, right_x, top, half_w, Inches(0.44), _NAVY)
        self._tx(slide, left_title,
                 LM + Inches(0.14), top + Emu(26000),
                 half_w - Inches(0.28), Inches(0.32),
                 bold=True, size=12, color=_WHITE)
        self._tx(slide, right_title,
                 right_x + Inches(0.14), top + Emu(26000),
                 half_w - Inches(0.28), Inches(0.32),
                 bold=True, size=12, color=_WHITE)

        bul_y = top + Inches(0.56)
        for i, (left_b, right_b) in enumerate(zip(
                (left_bullets or [])[:5],
                (right_bullets or [])[:5])):
            self._rect(slide, LM, bul_y + Emu(20000),
                       Emu(8000), Inches(0.45), _BLUE)
            self._tx(slide, left_b,
                     LM + Inches(0.14), bul_y,
                     half_w - Inches(0.18), Inches(0.72),
                     size=13, color=_SLATE, wrap=True)
            self._rect(slide, right_x, bul_y + Emu(20000),
                       Emu(8000), Inches(0.45), _NAVY)
            self._tx(slide, right_b,
                     right_x + Inches(0.14), bul_y,
                     half_w - Inches(0.18), Inches(0.72),
                     size=13, color=_SLATE, wrap=True)
            bul_y += Inches(0.82)

        # Handle unequal lists — render remaining items on either side
        for b in (left_bullets or [])[len(right_bullets or []):5]:
            self._rect(slide, LM, bul_y + Emu(20000), Emu(8000), Inches(0.45), _BLUE)
            self._tx(slide, b, LM + Inches(0.14), bul_y,
                     half_w - Inches(0.18), Inches(0.72), size=13, color=_SLATE, wrap=True)
            bul_y += Inches(0.82)
        for b in (right_bullets or [])[len(left_bullets or []):5]:
            self._rect(slide, right_x, bul_y + Emu(20000), Emu(8000), Inches(0.45), _NAVY)
            self._tx(slide, b, right_x + Inches(0.14), bul_y,
                     half_w - Inches(0.18), Inches(0.72), size=13, color=_SLATE, wrap=True)
            bul_y += Inches(0.82)
        return slide

    def add_data_chart(self, insight_headline: str, chart_path: str,
                       data_source: str = "") -> object:
        """Insight headline + embedded chart PNG."""
        slide = self._slide()
        self._header(slide, insight_headline)
        self._footer(slide)

        p = Path(chart_path)
        if p.exists():
            img_w = Inches(11.2)
            img_h = Inches(5.45)
            slide.shapes.add_picture(str(p),
                                     (W - img_w) / 2,
                                     HDR_H + Inches(0.12),
                                     img_w, img_h)
        else:
            self._tx(slide, "[ Chart visualization pending — data required ]",
                     LM, Inches(3.5), CW, Inches(0.55),
                     size=13, italic=True, color=_MUTED, align=PP_ALIGN.CENTER)

        if data_source:
            self._tx(slide, f"Source: {data_source}",
                     LM, Inches(6.62), CW, Inches(0.28),
                     size=8, italic=True, color=_MUTED)
        return slide

    def add_comparison_table(self, headline: str,
                             col_headers: list[str],
                             rows: list[list[str]]) -> object:
        """Feature/criteria comparison table."""
        slide = self._slide()
        self._header(slide, headline)
        self._footer(slide)

        n_rows = len(rows) + 1
        n_cols = len(col_headers) + 1
        tbl_h  = min(Inches(5.5), int(Pt(36) * n_rows))
        shape  = slide.shapes.add_table(
            n_rows, n_cols, int(LM), int(HDR_H + Inches(0.18)), int(CW), int(tbl_h))
        tbl = shape.table

        label_w = Inches(2.8)
        data_w  = int((CW - label_w) / max(len(col_headers), 1))
        tbl.columns[0].width = int(label_w)
        for i in range(1, n_cols):
            tbl.columns[i].width = data_w

        _cell(tbl.cell(0, 0), "", _NAVY, _WHITE, bold=True, size=11)
        for j, h in enumerate(col_headers):
            _cell(tbl.cell(0, j + 1), h, _NAVY, _WHITE, bold=True, size=11)

        for i, row in enumerate(rows):
            bg = _LIGHT if i % 2 == 0 else _WHITE
            _cell(tbl.cell(i + 1, 0), row[0] if row else "", bg, _NAVY, bold=True, size=11)
            for j, val in enumerate(row[1:], 1):
                _cell(tbl.cell(i + 1, j), str(val), bg, _SLATE, size=11)
        return slide

    def add_roadmap(self, headline: str, steps: list[dict],
                    chart_path: str = "") -> object:
        """Milestone roadmap — embeds PNG if available, else text boxes."""
        slide = self._slide()
        self._header(slide, headline)
        self._footer(slide)

        if chart_path and Path(chart_path).exists():
            img_w = Inches(12.2)
            img_h = Inches(5.0)
            slide.shapes.add_picture(chart_path,
                                     (W - img_w) / 2,
                                     HDR_H + Inches(0.15),
                                     img_w, img_h)
        else:
            n = min(len(steps), 6)
            if n:
                box_w = CW / n
                for i, s in enumerate(steps[:n]):
                    x = LM + box_w * i
                    self._rect(slide, x + Inches(0.04),
                               HDR_H + Inches(0.45),
                               box_w - Inches(0.08),
                               Inches(0.44), _BLUE)
                    self._tx(slide, s.get("title", f"Phase {i+1}"),
                             x + Inches(0.1), HDR_H + Inches(0.47),
                             box_w - Inches(0.2), Inches(0.38),
                             bold=True, size=12, color=_WHITE,
                             align=PP_ALIGN.CENTER)
                    if s.get("detail"):
                        self._tx(slide, s["detail"],
                                 x + Inches(0.1), HDR_H + Inches(1.05),
                                 box_w - Inches(0.2), Inches(1.5),
                                 size=11, color=_SLATE, wrap=True)
        return slide

    def add_recommendation(self, recommendation: str,
                           rationale: list[str]) -> object:
        """Framed recommendation box + numbered rationale."""
        slide = self._slide()
        self._header(slide, "Recommendation")
        self._footer(slide)

        # Warm recommendation box
        box_top = HDR_H + Inches(0.2)
        self._rect(slide, LM, box_top, CW, Inches(1.85), _WARM)
        self._tx(slide, "WE RECOMMEND:",
                 LM + Inches(0.22), box_top + Inches(0.1),
                 CW - Inches(0.44), Inches(0.26),
                 bold=True, size=8, color=_WHITE)
        self._tx(slide, recommendation,
                 LM + Inches(0.22), box_top + Inches(0.38),
                 CW - Inches(0.44), Inches(1.25),
                 bold=True, size=17, color=_WHITE, wrap=True)

        # Rationale
        rat_top = box_top + Inches(2.0)
        self._tx(slide, "RATIONALE",
                 LM, rat_top, CW, Inches(0.3),
                 bold=True, size=9, color=_NAVY)
        y = rat_top + Inches(0.35)
        for r in rationale[:4]:
            self._rect(slide, LM, y + Emu(22000), Emu(10000), Inches(0.44), _BLUE)
            self._tx(slide, r,
                     LM + Inches(0.18), y,
                     CW - Inches(0.22), Inches(0.65),
                     size=14, color=_SLATE, wrap=True)
            y += Inches(0.82)
        return slide

    def add_next_steps(self, steps: list[dict]) -> object:
        """Numbered action items with owner and timeline."""
        slide = self._slide()
        self._header(slide, "Next Steps")
        self._footer(slide)

        y = HDR_H + Inches(0.22)
        for i, s in enumerate(steps[:5]):
            action   = s.get("action", "")
            owner    = s.get("owner", "")
            timeline = s.get("timeline", "")

            # Number chip
            self._rect(slide, LM, y + Emu(12000),
                       Emu(50000), Emu(50000), _BLUE)
            self._tx(slide, f"{i + 1:02d}",
                     LM + Emu(6000), y + Emu(8000),
                     Emu(50000), Emu(50000),
                     bold=True, size=13, color=_WHITE,
                     align=PP_ALIGN.CENTER)

            self._tx(slide, action,
                     LM + Inches(0.62), y,
                     CW - Inches(4.0), Inches(0.56),
                     bold=True, size=15, color=_NAVY, wrap=True)

            meta = "  ·  ".join(x for x in [
                f"Owner: {owner}" if owner else "", timeline] if x)
            if meta:
                self._tx(slide, meta,
                         LM + Inches(0.62), y + Inches(0.55),
                         CW - Inches(4.0), Inches(0.3),
                         size=11, color=_MUTED)

            # Timeline chip on right
            if timeline:
                self._rect(slide, W - RM - Inches(3.0), y + Emu(12000),
                           Inches(3.0), Emu(50000), _LIGHT)
                self._tx(slide, timeline,
                         W - RM - Inches(2.9), y + Emu(8000),
                         Inches(2.8), Emu(50000),
                         size=11, color=_SLATE, align=PP_ALIGN.CENTER)
            y += Inches(1.12)
        return slide

    def add_appendix_cover(self) -> object:
        return self.add_section_divider("Appendix", "Supporting detail and data")

    # ── ISO Case Study slides ─────────────────────────────────────────────────

    def add_iso_cover(self, standard: str, clause: str, title: str,
                      date: str = "") -> object:
        """ISO case study cover — teal accent, standard badge."""
        slide = self._slide()
        self._rect(slide, 0, 0, W, H, _NAVY)
        self._rect(slide, 0, 0, Emu(42000), H, _TEAL)

        self._tx(slide, "LATITUDE MEDTECH LLC",
                 W - Inches(4.2), Inches(0.3), Inches(3.7), Inches(0.38),
                 size=10, bold=True, color=_TEAL, align=PP_ALIGN.RIGHT)
        self._tx(slide, _TAGLINE,
                 W - Inches(4.2), Inches(0.62), Inches(3.7), Inches(0.28),
                 size=9, italic=True, color=_MUTED, align=PP_ALIGN.RIGHT)

        # Standard badge
        self._rect(slide, Inches(0.85), Inches(1.7), Inches(4.0), Inches(0.52), _TEAL)
        self._tx(slide, standard,
                 Inches(0.98), Inches(1.76), Inches(3.8), Inches(0.42),
                 bold=True, size=15, color=_WHITE)

        # Clause label
        self._tx(slide, f"Clause {clause}",
                 Inches(0.85), Inches(2.42), Inches(6), Inches(0.4),
                 size=14, color=_MUTED)

        # Title
        self._tx(slide, title,
                 Inches(0.85), Inches(2.9), Inches(11.0), Inches(2.0),
                 bold=True, size=38, color=_WHITE)

        self._rect(slide, Inches(0.85), Inches(5.1), Inches(5), Emu(10000), _TEAL)

        if date:
            self._tx(slide, date,
                     Inches(0.85), Inches(5.35), Inches(5), Inches(0.3),
                     size=11, color=_MUTED)

        self._tx(slide, "ISO Coaching Series — Case Study",
                 Inches(0.85), Inches(6.9), Inches(6), Inches(0.35),
                 size=10, italic=True, color=_MUTED)
        self._tx(slide, "CONFIDENTIAL",
                 W - Inches(2.8), Inches(7.05), Inches(2.3), Inches(0.35),
                 size=9, color=_MUTED, align=PP_ALIGN.RIGHT)
        return slide

    def add_iso_overview(self, what_it_covers: str, why_it_matters: str,
                         key_requirement: str = "") -> object:
        """ISO clause overview — two-column brief."""
        slide = self._slide()
        self._header(slide, "Clause Overview", "What this clause covers and why it matters")
        self._footer(slide)

        half_w  = (CW - Inches(0.3)) / 2
        right_x = LM + half_w + Inches(0.3)
        top     = HDR_H + Inches(0.22)

        self._rect(slide, LM, top, half_w, Inches(0.4), _TEAL)
        self._tx(slide, "WHAT IT COVERS",
                 LM + Inches(0.14), top + Emu(22000),
                 half_w - Inches(0.28), Inches(0.28),
                 bold=True, size=9, color=_WHITE)
        self._tx(slide, what_it_covers,
                 LM, top + Inches(0.52), half_w, Inches(4.5),
                 size=14, color=_SLATE, wrap=True)

        self._rect(slide, right_x, top, half_w, Inches(0.4), _NAVY)
        self._tx(slide, "WHY IT MATTERS",
                 right_x + Inches(0.14), top + Emu(22000),
                 half_w - Inches(0.28), Inches(0.28),
                 bold=True, size=9, color=_WHITE)
        self._tx(slide, why_it_matters,
                 right_x, top + Inches(0.52), half_w, Inches(4.5),
                 size=14, color=_SLATE, wrap=True)

        if key_requirement:
            req_y = Inches(6.35)
            self._rect(slide, LM, req_y, CW, Inches(0.52), _LIGHT)
            self._tx(slide, f"Key requirement: {key_requirement}",
                     LM + Inches(0.2), req_y + Emu(16000), CW - Inches(0.4), Inches(0.38),
                     size=12, italic=True, color=_SLATE)
        return slide

    def add_iso_requirements(self, requirements: list[str]) -> object:
        """ISO requirements bullet list."""
        slide = self._slide()
        self._header(slide, "Regulatory Requirements")
        self._footer(slide)

        y = HDR_H + Inches(0.22)
        for req in requirements[:6]:
            self._rect(slide, LM, y + Emu(22000), Emu(10000), Inches(0.5), _TEAL)
            self._tx(slide, req,
                     LM + Inches(0.18), y,
                     CW - Inches(0.22), Inches(0.78),
                     size=14, color=_SLATE, wrap=True)
            y += Inches(0.92)
        return slide

    def add_iso_case_study_slide(self, title: str, company: str,
                                  scenario: str, finding: str,
                                  lesson: str, is_real: bool = True) -> object:
        """ISO case study — situation / finding / lesson layout."""
        slide = self._slide()
        badge_color = _NAVY if is_real else _TEAL
        badge_label = "REAL-WORLD CASE" if is_real else "HYPOTHETICAL SCENARIO"
        self._header(slide, title, badge_label)
        self._footer(slide)

        # Company badge
        self._rect(slide, LM, HDR_H + Inches(0.18), Inches(3.5), Inches(0.36), badge_color)
        self._tx(slide, company,
                 LM + Inches(0.12), HDR_H + Inches(0.22),
                 Inches(3.3), Inches(0.28),
                 bold=True, size=11, color=_WHITE)

        top = HDR_H + Inches(0.68)
        sections = [
            ("SCENARIO",  scenario, _BLUE),
            ("FINDING",   finding,  _WARM),
            ("LESSON",    lesson,   _TEAL),
        ]
        col_w = (CW - Inches(0.4)) / 3
        for i, (label, body, color) in enumerate(sections):
            x = LM + (col_w + Inches(0.2)) * i
            self._rect(slide, x, top, col_w, Inches(0.36), color)
            self._tx(slide, label,
                     x + Inches(0.12), top + Emu(18000),
                     col_w - Inches(0.24), Inches(0.26),
                     bold=True, size=8, color=_WHITE)
            self._tx(slide, body,
                     x, top + Inches(0.46),
                     col_w, Inches(4.5),
                     size=13, color=_SLATE, wrap=True)
        return slide

    def add_iso_mistakes(self, mistakes: list[str]) -> object:
        """Common implementation mistakes with warning markers."""
        slide = self._slide()
        self._header(slide, "Common Mistakes", "What practitioners get wrong — and how to avoid it")
        self._footer(slide)

        y = HDR_H + Inches(0.22)
        for m in mistakes[:5]:
            self._rect(slide, LM, y + Emu(20000), Emu(10000), Inches(0.5), _RED)
            self._tx(slide, m,
                     LM + Inches(0.18), y,
                     CW - Inches(0.22), Inches(0.8),
                     size=14, color=_SLATE, wrap=True)
            y += Inches(1.0)
        return slide

    def add_iso_key_terms(self, terms: dict[str, str],
                          career_note: str = "") -> object:
        """Key terms glossary + optional career connection."""
        slide = self._slide()
        self._header(slide, "Key Terms & Career Connection")
        self._footer(slide)

        # Terms (left ~70%)
        term_w = CW * 0.68
        sep_x  = LM + term_w + Inches(0.25)
        top    = HDR_H + Inches(0.2)

        self._tx(slide, "KEY TERMS",
                 LM, top, term_w, Inches(0.28),
                 bold=True, size=9, color=_NAVY)
        y = top + Inches(0.35)
        for term, defn in list(terms.items())[:6]:
            self._rect(slide, LM, y + Emu(14000), Emu(8000), Inches(0.42), _TEAL)
            self._tx(slide, f"{term}: {defn}",
                     LM + Inches(0.14), y,
                     term_w - Inches(0.18), Inches(0.72),
                     size=12, color=_SLATE, wrap=True)
            y += Inches(0.78)

        # Career note (right ~30%)
        if career_note:
            right_w = W - sep_x - RM
            self._rect(slide, sep_x, top, Emu(7000), Inches(5.5), _RULE)
            self._rect(slide, sep_x + Inches(0.25), top, right_w, Inches(0.36), _NAVY)
            self._tx(slide, "CAREER CONNECTION",
                     sep_x + Inches(0.37), top + Emu(18000),
                     right_w - Inches(0.24), Inches(0.26),
                     bold=True, size=8, color=_WHITE)
            self._tx(slide, career_note,
                     sep_x + Inches(0.25), top + Inches(0.46),
                     right_w, Inches(5.0),
                     size=12, color=_SLATE, wrap=True)
        return slide


# ── Table cell helper ─────────────────────────────────────────────────────────

def _cell(cell, text: str, bg: RGBColor, fg: RGBColor, *,
          bold: bool = False, size: int = 11):
    cell.fill.solid()
    cell.fill.fore_color.rgb = bg
    tf = cell.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.space_before = Pt(4)
    para.space_after  = Pt(4)
    if para.runs:
        run = para.runs[0]
    else:
        run = para.add_run()
    run.text           = text
    run.font.bold      = bold
    run.font.size      = Pt(size)
    run.font.color.rgb = fg
    run.font.name      = FONT
