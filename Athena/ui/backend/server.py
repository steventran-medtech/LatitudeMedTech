"""
Latitude MedTech — Local Backend API
======================================
FastAPI server that powers the desktop UI.
Runs on localhost:8000.

Start: python server.py
"""

import os
import sys
import re
import json
import secrets
import subprocess
import asyncio
import logging
import threading
from logging.handlers import RotatingFileHandler
from pathlib import Path
from datetime import datetime, timedelta
from typing import Literal, Optional

from fastapi import FastAPI, WebSocket, WebSocketDisconnect, BackgroundTasks, Request, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, JSONResponse
from starlette.middleware.base import BaseHTTPMiddleware
from slowapi import Limiter, _rate_limit_exceeded_handler
from slowapi.util import get_remote_address
from slowapi.errors import RateLimitExceeded
from pydantic import BaseModel
from dotenv import load_dotenv

ATHENA      = Path(__file__).resolve().parents[2]        # ui/backend/ -> ui/ -> Athena/
AGENTS_DIR  = ATHENA / "agents"
VENV_PYTHON = ATHENA / "voice" / "venv" / "Scripts" / "python.exe"

# ── App version (single source of truth: Athena/VERSION.json) ──────────────────
def _load_version_info() -> dict:
    """Read the canonical version metadata. Never raises — falls back to a
    dev placeholder so a missing/corrupt VERSION.json can't block startup."""
    try:
        return json.loads((ATHENA / "VERSION.json").read_text(encoding="utf-8"))
    except Exception:
        return {"version": "0.0.0", "released": "unknown", "channel": "dev"}

VERSION_INFO = _load_version_info()
APP_VERSION  = VERSION_INFO.get("version", "0.0.0")

load_dotenv(ATHENA / "voice" / ".env")

sys.path.insert(0, str(AGENTS_DIR))

try:
    from memory import Memory
    mem = Memory()
except ImportError:
    mem = None

# LangGraph coaching orchestrator (Phase 1A). Loaded lazily so a missing
# langgraph install never blocks server startup.
_orch = None
def _orchestrator():
    global _orch
    if _orch is None:
        import orchestrator as _orch  # noqa: F811
    return _orch

# ── Session token (generated once per server start) ───────────────────────────
_KEY_FILE = ATHENA / "voice" / ".athena.key"

def _get_or_create_token() -> str:
    if _KEY_FILE.exists():
        token = _KEY_FILE.read_text().strip()
        if len(token) >= 32:
            return token
    token = secrets.token_hex(32)
    _KEY_FILE.write_text(token)
    try:
        import stat
        _KEY_FILE.chmod(stat.S_IRUSR | stat.S_IWUSR)
    except Exception:
        pass
    return token

SESSION_TOKEN = _get_or_create_token()

# ── Rate limiter ──────────────────────────────────────────────────────────────
limiter = Limiter(key_func=get_remote_address, default_limits=["120/minute"])

# ── Security helpers ──────────────────────────────────────────────────────────

_SAFE_FILENAME = re.compile(r'^[\w\-. ]+$')
_SAFE_ARG      = re.compile(r'^[\w\-. /:\[\]@#]+$')

_MAX_BODY_BYTES = 10 * 1024 * 1024  # 10 MB hard limit on request bodies

def _safe_filename(name: str) -> str:
    """Reject filenames with path traversal characters."""
    if not name or not _SAFE_FILENAME.match(name) or '..' in name or '/' in name or '\\' in name:
        raise HTTPException(status_code=400, detail="Invalid filename")
    return name

def _safe_arg(arg: str, max_len: int = 500) -> str:
    """Sanitise a free-text override arg — strip shell metacharacters."""
    if not arg:
        return ""
    sanitised = re.sub(r'[;&|`$<>{}]', '', arg)[:max_len]
    return sanitised.strip()

def _extract_md_summary(path: "Path", max_chars: int = 220) -> str:
    """Return the first substantive text line from a markdown file (skips frontmatter and headings)."""
    try:
        text = path.read_text(encoding="utf-8", errors="replace")
        if text.startswith("---"):
            end = text.find("\n---", 3)
            if end != -1:
                text = text[end + 4:].strip()
        for line in text.splitlines():
            stripped = line.strip()
            if stripped and not stripped.startswith("#") and not stripped.startswith("---"):
                return stripped[:max_chars]
    except Exception:
        pass
    return ""

# ── Audit logging (SOC II CC7.2) ──────────────────────────────────────────────
# Rotating at 5 MB, keeping 5 backups → max 30 MB on-disk; never silently drops.

_AUDIT_LOG = ATHENA / "logs" / "audit.log"
_AUDIT_LOG.parent.mkdir(parents=True, exist_ok=True)

_audit_logger = logging.getLogger("latitude.audit")
_audit_logger.setLevel(logging.INFO)
_audit_logger.propagate = False
_audit_handler = RotatingFileHandler(
    str(_AUDIT_LOG), maxBytes=5 * 1024 * 1024, backupCount=5, encoding="utf-8"
)
_audit_handler.setFormatter(logging.Formatter("%(message)s"))
_audit_logger.addHandler(_audit_handler)

def _audit(action: str, detail: str = "", request: "Request | None" = None) -> None:
    """Append a one-line audit record (SOC II CC7.2). Rotates at 5 MB, 5 backups."""
    ts     = datetime.now().isoformat(timespec="seconds")
    origin = request.headers.get("origin", "local") if request else "system"
    _audit_logger.info(f"{ts}  {action:<30}  origin={origin}  {detail}")

# ── Database file protection (SOC II CC6.1) ───────────────────────────────────

def _harden_db_permissions() -> None:
    """Set restrictive file permissions on the SQLite database and key file."""
    import stat
    targets = [
        ATHENA / "memory" / "latitude_memory.db",
        _KEY_FILE,
    ]
    for p in targets:
        if p.exists():
            try:
                p.chmod(stat.S_IRUSR | stat.S_IWUSR)
            except Exception:
                pass

# ── Session-token authentication middleware (SOC II CC6.6) ────────────────────
# All API requests must carry X-Athena-Key: <SESSION_TOKEN> (or ?token= for
# WebSocket handshakes).  Exempt: / · /api/version · /api/auth/token.

_AUTH_EXEMPT = frozenset({"/", "/api/version", "/api/auth/token"})

class AuthMiddleware(BaseHTTPMiddleware):
    async def dispatch(self, request: Request, call_next):
        path = request.url.path
        if path in _AUTH_EXEMPT or path.startswith("/ws"):
            return await call_next(request)
        key = (request.headers.get("X-Athena-Key", "")
               or request.query_params.get("token", ""))
        if not secrets.compare_digest(key.encode(), SESSION_TOKEN.encode()):
            _audit("auth_failure",
                   f"path={path} ip={getattr(request.client, 'host', 'unknown')}",
                   request)
            return JSONResponse(status_code=401, content={"error": "Unauthorized"})
        return await call_next(request)

# ── Request body size limit middleware ────────────────────────────────────────
class BodySizeLimitMiddleware(BaseHTTPMiddleware):
    async def dispatch(self, request, call_next):
        content_length = request.headers.get("content-length")
        if content_length and int(content_length) > _MAX_BODY_BYTES:
            return JSONResponse(status_code=413, content={"error": "Request body too large"})
        return await call_next(request)

# ── Security headers middleware ───────────────────────────────────────────────
class SecurityHeadersMiddleware(BaseHTTPMiddleware):
    async def dispatch(self, request, call_next):
        response = await call_next(request)
        response.headers["X-Content-Type-Options"]  = "nosniff"
        response.headers["X-Frame-Options"]         = "DENY"
        response.headers["Referrer-Policy"]         = "strict-origin-when-cross-origin"
        response.headers["X-XSS-Protection"]        = "1; mode=block"
        # Tight CSP for local-only app
        response.headers["Content-Security-Policy"] = (
            "default-src 'self'; "
            "script-src 'self' 'unsafe-inline' 'unsafe-eval'; "
            "style-src 'self' 'unsafe-inline' https://fonts.googleapis.com; "
            "font-src 'self' https://fonts.gstatic.com; "
            "connect-src 'self' ws://localhost:* http://localhost:*; "
            "img-src 'self' data: blob:;"
        )
        return response

app = FastAPI(title="Latitude MedTech", version=APP_VERSION,
              docs_url=None, redoc_url=None)  # disable public docs
app.state.limiter = limiter
app.add_exception_handler(RateLimitExceeded, _rate_limit_exceeded_handler)
app.add_middleware(SecurityHeadersMiddleware)
app.add_middleware(BodySizeLimitMiddleware)
app.add_middleware(AuthMiddleware)           # SOC II CC6.6: token auth on every request

# Harden file permissions on startup (SOC II CC6.1)
_harden_db_permissions()
_audit("server_start", f"version={APP_VERSION}")

# ── Voice bridge ──────────────────────────────────────────────────────────────
sys.path.insert(0, str(ATHENA / "voice"))
try:
    from voice_bridge import router as voice_router, voice_websocket_endpoint, _notification_queue as _voice_queue
    app.include_router(voice_router)
    VOICE_AVAILABLE = True
    print("[voice] Voice bridge loaded — OWW + Whisper + Kokoro preloading in background")
except Exception as _ve:
    VOICE_AVAILABLE = False
    _voice_queue = None
    print(f"[voice] Bridge not loaded: {_ve}")

# Track review-queue IDs already spoken so the batch announcement never repeats
# an item that was announced in an earlier agent-completion cycle this session.
_announced_review_ids: set = set()


def _vary_voice_phrase(base: str) -> str:
    """Return a naturally varied rephrasing of `base` via Claude Haiku for voice delivery.
    Preserves all titles and names exactly.  Falls back to `base` on any error."""
    try:
        import anthropic as _ant
        _key = os.getenv("ANTHROPIC_API_KEY", "")
        if not _key:
            return base
        _client = _ant.Anthropic(api_key=_key)
        _resp = _client.messages.create(
            model="claude-haiku-4-5",
            max_tokens=80,
            system=(
                "You are Athena, a warm and professional British English voice assistant "
                "for Latitude MedTech LLC. Rephrase the given phrase as a natural spoken "
                "notification — one sentence, no markdown. "
                "Preserve all specific titles, names, and numbers exactly as given."
            ),
            messages=[{"role": "user", "content": base}],
        )
        return _resp.content[0].text.strip()
    except Exception:
        return base


app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:3000", "http://127.0.0.1:3000",
                   "app://.", "file://"],           # Electron origins
    allow_methods=["GET", "POST", "DELETE"],
    allow_headers=["Content-Type", "Authorization", "X-Athena-Key"],
    allow_credentials=False,
)

# ── WebSocket manager ─────────────────────────────────────────────────────────

class ConnectionManager:
    def __init__(self):
        self.active: list[WebSocket] = []

    async def connect(self, ws: WebSocket):
        await ws.accept()
        self.active.append(ws)

    def disconnect(self, ws: WebSocket):
        self.active.remove(ws)

    async def broadcast(self, data: dict):
        dead = []
        for ws in self.active:
            try:
                await ws.send_json(data)
            except Exception:
                dead.append(ws)
        for ws in dead:
            self.active.remove(ws)

manager = ConnectionManager()
running_agents = set()  # Prevent duplicate runs

_shutting_down = False          # set True once /api/shutdown is called
_goodbye_task: "asyncio.Task | None" = None  # delayed goodbye after WS disconnect


async def _ws_goodbye():
    """Plays TTS farewell after all WebSocket clients disconnect.

    Waits 3 s for a page-refresh reconnect before acting.  Guarded by
    _shutting_down so the in-app exit button flow never double-speaks.
    """
    await asyncio.sleep(3)
    if manager.active or _shutting_down:
        return
    stop_script = ATHENA / "ui" / "stop_athena.ps1"
    phrase = _random.choice(_GOODBYES)
    await asyncio.get_event_loop().run_in_executor(None, _speak_phrase, phrase)
    try:
        flags = (getattr(subprocess, "CREATE_NO_WINDOW", 0) |
                 getattr(subprocess, "CREATE_NEW_PROCESS_GROUP", 0))
        subprocess.Popen(
            ["powershell.exe", "-ExecutionPolicy", "Bypass", "-NoProfile",
             "-File", str(stop_script)],
            creationflags=flags, close_fds=True,
        )
    except Exception:
        pass
    import os as _os; _os._exit(0)


# ── Agent runner ──────────────────────────────────────────────────────────────

async def run_agent(agent_name: str, script: str, args: list = None, context: str = ""):
    """Run an agent script and stream output via WebSocket."""
    if agent_name in running_agents:
        await manager.broadcast({"type": "agent_log", "agent": agent_name, "line": f"[SKIPPED] {agent_name} already running"})
        return
    running_agents.add(agent_name)
    _audit("agent_start", f"agent={agent_name} args={args or []}")
    msg = {"type": "agent_start", "agent": agent_name, "ts": datetime.now().isoformat()}
    if context:
        msg["context"] = context
    await manager.broadcast(msg)

    # Snapshot pending review IDs before running so we can detect new items added by this agent.
    run_started_at = datetime.now().isoformat()

    cmd = [str(VENV_PYTHON), str(AGENTS_DIR / script)] + (args or [])
    proc = await asyncio.create_subprocess_exec(
        *cmd,
        stdout=asyncio.subprocess.PIPE,
        stderr=asyncio.subprocess.STDOUT,
        cwd=str(AGENTS_DIR),
    )

    lines = []
    async for line in proc.stdout:
        text = line.decode('utf-8', errors='replace').rstrip()
        lines.append(text)
        await manager.broadcast({"type": "agent_log", "agent": agent_name, "line": text})

    await proc.wait()
    running_agents.discard(agent_name)
    status = "success" if proc.returncode == 0 else "error"
    _audit("agent_done", f"agent={agent_name} status={status} rc={proc.returncode}")

    # Parse the last few stdout lines for a structured JSON result (deck_agent outputs a
    # JSON summary as its final line: {status, path, filename, title, slide_count, deck_type})
    result_extras = {}
    for line in reversed(lines[-8:]):
        try:
            parsed = json.loads(line.strip())
            if isinstance(parsed, dict) and "status" in parsed:
                result_extras = {k: parsed[k] for k in
                                 ("path", "filename", "title", "slide_count", "deck_type")
                                 if k in parsed}
                break
        except (json.JSONDecodeError, ValueError):
            continue

    # Detect review deliverables generated by this agent run.
    review_added = 0
    new_items: list = []
    if status == "success" and mem:
        try:
            pending   = mem.get_pending_reviews()
            new_items = [p for p in pending if p.get("timestamp", "") >= run_started_at]
            review_added = len(new_items)
        except Exception:
            pass

    await manager.broadcast({"type": "agent_done", "agent": agent_name, "status": status,
                              "ts": datetime.now().isoformat(),
                              "review_added": review_added, **result_extras})

    # Voice batch-announce newly queued review items (SOC II CC6.6 — user awareness).
    if new_items and _voice_queue is not None:
        try:
            n = len(new_items)
            if n == 1:
                base = f"Your {new_items[0]['title']} deliverables are ready for your review."
            else:
                listed = ", ".join(p["title"] for p in new_items[:3])
                tail   = f", and {n - 3} more" if n > 3 else ""
                base   = f"{n} deliverables ready for your review: {listed}{tail}."
            _voice_queue.append(_vary_voice_phrase(base))
            _announced_review_ids.update(p["id"] for p in new_items)
        except Exception:
            pass

    return {"status": status, "lines": lines}


# ── Routes ────────────────────────────────────────────────────────────────────

@app.get("/")
def root():
    return {"status": "ok", "service": "Latitude MedTech API", "version": APP_VERSION}


@app.get("/api/version")
def get_version():
    """Current Athena version plus the full changelog (raw Markdown).

    Source of truth is Athena/VERSION.json (cached at startup as VERSION_INFO);
    the changelog is read fresh so edits show up without a server restart.
    """
    changelog = ""
    try:
        changelog = (ATHENA / "CHANGELOG.md").read_text(encoding="utf-8")
    except Exception:
        pass
    return {**VERSION_INFO, "changelog": changelog}


@app.get("/api/auth/token")
def get_token():
    """Returns the session token for the local UI. No auth required (localhost only)."""
    return {"token": SESSION_TOKEN}


_dashboard_cache = {"data": None, "ts": 0}

@app.get("/api/dashboard")
def get_dashboard():
    global _dashboard_cache
    # Cache for 60 seconds
    if _dashboard_cache["data"] and (datetime.now().timestamp() - _dashboard_cache["ts"]) < 60:
        return _dashboard_cache["data"]
    if not mem:
        return {"error": "Memory not available"}
    try:
        report = mem.get_token_report(days=30)
        kb     = mem.get_kb_stats()
        topics = mem.get_recent_topics(days=90)
        result = {
            "token_report": report,
            "kb_stats":     kb,
            "recent_topics":topics[:20],
            "generated_at": datetime.now().isoformat(),
        }
        _dashboard_cache = {"data": result, "ts": datetime.now().timestamp()}
        return result
    except Exception as e:
        return {"error": str(e), "token_report": {}, "kb_stats": {}, "recent_topics": []}


# ── Voice session log ─────────────────────────────────────────────────────────

SESSIONS_FILE = ATHENA / "voice" / "sessions.jsonl"

class SessionLogBody(BaseModel):
    session_id: str
    started_at: str
    ended_at: str
    duration_secs: int
    queries: int
    device: Optional[str] = ""
    is_new: bool = True

class SessionStartBody(BaseModel):
    session_id: str
    started_at: str
    is_new: bool = True
    device: Optional[str] = ""

@app.post("/api/sessions/start")
async def session_start(body: SessionStartBody, request: Request):
    """Log the moment a voice session begins so partial/crashed sessions are visible."""
    record = {**body.dict(), "status": "started"}
    with open(SESSIONS_FILE, "a", encoding="utf-8") as f:
        f.write(json.dumps(record) + "\n")
    _audit("session_start", f"id={body.session_id} new={body.is_new}", request)
    return {"status": "logged"}

@app.post("/api/sessions/log")
async def log_session(body: SessionLogBody, request: Request):
    record = {**body.dict(), "status": "ended"}
    with open(SESSIONS_FILE, "a", encoding="utf-8") as f:
        f.write(json.dumps(record) + "\n")
    _audit("session_end", f"id={body.session_id} dur={body.duration_secs}s q={body.queries}", request)
    return {"status": "logged"}

@app.get("/api/sessions")
def get_sessions(limit: int = 20):
    if not SESSIONS_FILE.exists():
        return {"sessions": []}
    sessions = []
    with open(SESSIONS_FILE, encoding="utf-8") as f:
        for line in f:
            line = line.strip()
            if line:
                try:
                    sessions.append(json.loads(line))
                except Exception:
                    pass
    return {"sessions": list(reversed(sessions[-limit:]))}


@app.get("/api/drafts")
def list_drafts():
    drafts_dir = ATHENA / 'content' / 'drafts'
    if not drafts_dir.exists():
        return {"drafts": []}
    files = sorted(drafts_dir.glob('*.md'), key=lambda f: f.stat().st_mtime, reverse=True)
    drafts = []
    for f in files[:20]:
        content = f.read_text(encoding='utf-8', errors='replace')
        title   = f.stem
        # Extract title from frontmatter
        for line in content.splitlines():
            if line.startswith('title:'):
                title = line.replace('title:', '').strip()
                break
        drafts.append({
            "filename": f.name,
            "title":    title,
            "modified": datetime.fromtimestamp(f.stat().st_mtime).isoformat(),
            "size":     f.stat().st_size,
            "path":     str(f),
        })
    return {"drafts": drafts}


@app.get("/api/drafts/{filename}")
def get_draft(filename: str):
    drafts_dir = ATHENA / 'content' / 'drafts'
    f = drafts_dir / filename
    if not f.exists():
        return JSONResponse(status_code=404, content={"error": "Not found"})
    return {"filename": filename, "content": f.read_text(encoding='utf-8', errors='replace')}


@app.get("/api/briefings")
def list_briefings():
    briefings_dir = ATHENA / 'briefings'
    if not briefings_dir.exists():
        return {"briefings": []}
    files = sorted(
        [f for f in briefings_dir.glob('*.md') if not f.name.startswith('.')],
        key=lambda f: f.stat().st_mtime, reverse=True
    )
    briefings = []
    for f in files[:10]:
        title, fm_date = _briefing_meta(f)
        # Prefer a date-prefixed filename (YYYY-MM-DD…); else the frontmatter date.
        name_date = f.name[:10] if f.name[:4].isdigit() and f.name[4] == '-' else ""
        briefings.append({
            "filename": f.name,
            "date":     name_date or fm_date or f.name[:10],
            "title":    title,
            "modified": datetime.fromtimestamp(f.stat().st_mtime).isoformat(),
            "path":     str(f),
        })
    return {"briefings": briefings}


def _briefing_meta(f):
    """Return (title, date). Title: frontmatter `title:`, else first H1, else prettified
    filename. Date: frontmatter `date:` if present, else ""."""
    title, fm_date = "", ""
    try:
        with open(f, 'r', encoding='utf-8') as fh:
            in_front = False
            for i, line in enumerate(fh):
                if i > 60:
                    break
                s = line.strip()
                if i == 0 and s == '---':
                    in_front = True
                    continue
                if in_front:
                    if s == '---':
                        in_front = False
                        continue
                    low = s.lower()
                    if not title and low.startswith('title:'):
                        title = s.split(':', 1)[1].strip().strip('"\'')
                    elif not fm_date and low.startswith('date:'):
                        fm_date = s.split(':', 1)[1].strip().strip('"\'')
                elif not title and s.startswith('# '):
                    title = s[2:].strip()
    except OSError:
        pass
    if not title:
        title = f.stem.replace('_', ' ').strip().title()
    return title, fm_date


@app.get("/api/briefings/{filename}")
def get_briefing(filename: str):
    briefings_dir = ATHENA / 'briefings'
    f = briefings_dir / filename
    if not f.exists():
        return JSONResponse(status_code=404, content={"error": "Not found"})
    return {"filename": filename, "content": f.read_text(encoding='utf-8', errors='replace')}


@app.get("/api/kb/stats")
def kb_stats():
    if not mem:
        return {"error": "Memory not available"}
    return mem.get_kb_stats()


# ── HR + Learning routes ──────────────────────────────────────────────────────

# Canonical agent roster (label + tier) so the Workforce view always renders a
# full row set even before the first HR Review has run.
WORKFORCE_ROSTER = [
    ("content",         "Content Agent",         "Manager"),
    ("briefing",        "Briefing Agent",        "Senior Associate"),
    ("iso",             "ISO Coach Agent",       "Manager"),
    ("coaching",        "Coaching Brief Agent",  "Manager"),
    ("fda",             "FDA Agent",             "Manager"),
    ("eu_mdr",          "EU MDR Agent",          "Manager"),
    ("rag",             "RAG Ingestion Agent",   "Senior Associate"),
    ("consulting",      "Consulting Agent",      "Senior Manager"),
    ("ma_intelligence", "M&A Intelligence Agent","Senior Manager"),
    ("hr",              "HR / L&D Manager Agent","Business Function"),
    ("voice_bridge",    "Voice Assistant",       "Associate"),
]


@app.get("/api/hr/health")
def hr_health():
    if not mem:
        return {"error": "Memory not available"}
    # Stored flag/error data keyed by agent (populated by the last HR Review).
    stored = {a["agent"]: a for a in mem.get_all_agent_health()}
    agents = []
    for key, label, tier in WORKFORCE_ROSTER:
        row = dict(stored.get(key, {"agent": key, "flag_status": "green",
                                    "error_count_7d": 0, "learning_7d": 0,
                                    "flag_reason": None}))
        row.setdefault("agent", key)
        row["label"] = label
        row["tier"]  = tier
        # Always overlay live timestamps so they never read "Never" while
        # activity exists in the DB under any name variant.
        live_run   = mem.get_agent_last_run(key)
        live_learn = mem.get_last_learning(key)
        if live_run:
            row["last_run"] = live_run
        if live_learn:
            row["last_learning"] = live_learn["timestamp"]
        agents.append(row)
    # Include any stored agents not in the roster (forward-compat).
    for key, a in stored.items():
        if key not in {k for k, _, _ in WORKFORCE_ROSTER}:
            agents.append(a)
    return {"agents": agents}


@app.get("/api/hr/skills")
def hr_skills():
    """Live skill/KB accumulation per agent, queried directly from the memory DB."""
    if not mem:
        return {"error": "Memory not available"}
    result = {}
    for key, label, tier in WORKFORCE_ROSTER:
        acc = mem.get_skill_accumulation(key)
        result[key] = {
            "label": label, "tier": tier,
            "total_items":  acc["total_items"],
            "total_chunks": acc["total_chunks"],
            "domains":      acc["domains"],
            "last":         acc["last"],
        }
    return {"skills": result}


@app.post("/api/agents/skills-profile")
async def trigger_skills_profile(background_tasks: BackgroundTasks):
    """Regenerate all agent skill/KB profile .md files."""
    background_tasks.add_task(run_agent, "skills_profile", "skills_profile.py")
    return {"status": "started", "agent": "skills_profile"}


@app.get("/api/hr/learning")
def hr_learning(days: int = 7):
    if not mem:
        return {"error": "Memory not available"}
    return {"stats": mem.get_learning_stats(days=days)}


@app.post("/api/agents/hr")
async def trigger_hr(background_tasks: BackgroundTasks):
    # Block HR review if any learning agent is currently running
    if "agent_learning" in running_agents:
        return JSONResponse(status_code=409,
            content={"error": "Learning in progress — wait for it to complete before running HR Review",
                     "blocked_by": "agent_learning"})
    # Also block if HR is already running (prevents double-click double-run)
    if "hr_agent" in running_agents:
        return JSONResponse(status_code=409,
            content={"error": "HR Review is already running", "blocked_by": "hr_agent"})
    background_tasks.add_task(run_agent, "hr_agent", "hr_agent.py")
    return {"status": "started", "agent": "hr_agent"}


@app.post("/api/agents/learn")
async def trigger_learning(request: Request, background_tasks: BackgroundTasks):
    body   = await request.json()
    agent  = body.get("agent", "")
    args   = ["--agent", agent] if agent else []
    background_tasks.add_task(run_agent, "agent_learning", "agent_learning.py", args)
    return {"status": "started", "agent": "agent_learning", "target": agent or "all"}


@app.post("/api/agents/run-all")
async def trigger_run_all(background_tasks: BackgroundTasks):
    """Run all standard agents in parallel (briefing, content, RAG, consulting, M&A, ISO)."""
    batch = [
        ("briefing_agent",        "briefing_agent.py",        []),
        ("content_agent",         "content_agent.py",         []),
        ("rag_agent",             "rag_agent.py",             []),
        ("consulting_agent",      "consulting_agent.py",      []),
        ("ma_intelligence_agent", "ma_intelligence_agent.py", []),
        ("iso_coach",             "iso_coach_agent.py",       ["--next"]),
    ]
    for name, script, args in batch:
        background_tasks.add_task(run_agent, name, script, args)
    return {"status": "started", "agents": [b[0] for b in batch]}


@app.post("/api/agents/consulting")
async def trigger_consulting(request: Request, background_tasks: BackgroundTasks):
    body = await request.json()
    mode = body.get("mode", "learn")   # "learn" | "report"
    override = _safe_arg(body.get("override", ""))
    args = ["--generate", "report"] if mode == "report" else []
    context = override or (mode if mode != "learn" else "")
    background_tasks.add_task(run_agent, "consulting_agent", "consulting_agent.py", args, context)
    return {"status": "started", "agent": "consulting_agent", "mode": mode}


@app.post("/api/agents/marketing")
async def trigger_marketing(request: Request, background_tasks: BackgroundTasks):
    body   = await request.json()
    mode   = body.get("mode", "brief")   # "brief" | "plan" | "events" | "scorecard" | "learn"
    target = body.get("target", "")
    if mode == "plan":
        args = ["--plan"]
    elif mode == "events":
        args = ["--events"]
    elif mode == "scorecard":
        args = ["--scorecard"]
    elif mode == "learn":
        args = ["learn"]
    elif mode == "outreach" and target:
        args = ["--outreach", _safe_arg(target)]
    else:
        args = []  # default: brief
    background_tasks.add_task(run_agent, "marketing_agent", "marketing_agent.py", args)
    return {"status": "started", "agent": "marketing_agent", "mode": mode}


@app.get("/api/marketing/pipeline")
def get_marketing_pipeline():
    """Return pipeline summary counts and high-priority next actions."""
    import sqlite3 as _sq
    db = ATHENA / "ops" / "marketing" / "pipeline.db"
    if not db.exists():
        return {"summary": [], "actions": [], "kpis": {}}
    with _sq.connect(db) as conn:
        conn.row_factory = _sq.Row
        summary = [dict(r) for r in conn.execute(
            "SELECT type, status, COUNT(*) as count FROM targets GROUP BY type, status ORDER BY type, status"
        ).fetchall()]
        actions = [dict(r) for r in conn.execute(
            "SELECT name, type, org, status, next_action, next_date FROM targets "
            "WHERE priority = 1 AND status NOT IN ('converted','closed') "
            "ORDER BY CASE WHEN next_date IS NULL THEN '9999' ELSE next_date END LIMIT 10"
        ).fetchall()]
        row = conn.execute(
            "SELECT COUNT(*) as total, "
            "SUM(CASE WHEN status!='identified' THEN 1 ELSE 0 END) as contacted, "
            "SUM(CASE WHEN status='in_discussion' THEN 1 ELSE 0 END) as meetings, "
            "SUM(CASE WHEN status='converted' THEN 1 ELSE 0 END) as converted "
            "FROM targets"
        ).fetchone()
        kpis = dict(row) if row else {}
    return {"summary": summary, "actions": actions, "kpis": kpis}


@app.post("/api/marketing/pipeline/update")
async def update_pipeline_target(request: Request):
    """Update a target's status from the UI."""
    import sqlite3 as _sq
    body   = await request.json()
    name   = _safe_arg(body.get("name", ""))
    status = _safe_arg(body.get("status", ""))
    note   = _safe_arg(body.get("note", ""))
    VALID  = {"identified", "contacted", "in_discussion", "converted", "closed"}
    if not name or status not in VALID:
        raise HTTPException(status_code=400, detail="Invalid name or status")
    db = ATHENA / "ops" / "marketing" / "pipeline.db"
    if not db.exists():
        raise HTTPException(status_code=404)
    with _sq.connect(db) as conn:
        conn.execute(
            "UPDATE targets SET status=?, notes=COALESCE(notes,'')||'\n'||?, updated_at=date('now') WHERE name LIKE ?",
            (status, f"[{datetime.now().strftime('%Y-%m-%d')}] {note}", f"%{name}%")
        )
        conn.commit()
    _audit("pipeline_update", f"target={name} status={status}", request)
    return {"status": "updated", "target": name}


@app.get("/api/marketing/outputs")
def list_marketing_outputs():
    """List all generated marketing output files."""
    mdir = ATHENA / "ops" / "marketing"
    if not mdir.exists():
        return {"files": []}
    files = []
    for f in sorted(mdir.glob("*.md"), key=lambda x: x.stat().st_mtime, reverse=True):
        files.append({"filename": f.name, "label": f.stem.replace("_", " ").title(),
                      "modified": datetime.fromtimestamp(f.stat().st_mtime).strftime("%Y-%m-%d")})
    out_dir = mdir / "outreach"
    if out_dir.exists():
        for f in sorted(out_dir.glob("*.md"), key=lambda x: x.stat().st_mtime, reverse=True):
            files.append({"filename": f"outreach/{f.name}",
                          "label": f.stem.replace("_outreach_", " → ").replace("_", " ").title(),
                          "modified": datetime.fromtimestamp(f.stat().st_mtime).strftime("%Y-%m-%d")})
    return {"files": files}


@app.get("/api/marketing/outputs/{filename:path}")
def get_marketing_output(filename: str):
    """Read a specific marketing output file."""
    parts = Path(filename).parts
    if ".." in parts or len(parts) > 2:
        raise HTTPException(status_code=400, detail="Invalid path")
    path = ATHENA / "ops" / "marketing" / filename
    if not path.exists() or path.suffix != ".md":
        raise HTTPException(status_code=404)
    return {"content": path.read_text(encoding="utf-8"), "filename": filename}


@app.post("/api/agents/ma")
async def trigger_ma(request: Request, background_tasks: BackgroundTasks):
    body  = await request.json()
    mode  = body.get("mode", "learn")   # "learn" | "analyse" | "historical"
    topic = _safe_arg(body.get("topic", ""))
    if mode == "analyse":
        args = ["--analyse"] + (["--topic", topic] if topic else [])
    elif mode == "historical":
        args = ["--historical"]
    else:
        args = []
    context = topic or (mode if mode != "learn" else "")
    background_tasks.add_task(run_agent, "ma_intelligence_agent", "ma_intelligence_agent.py", args, context)
    return {"status": "started", "agent": "ma_intelligence_agent", "mode": mode}


# ── Agent triggers ────────────────────────────────────────────────────────────

@app.post("/api/agents/rag")
async def trigger_rag(request: Request, background_tasks: BackgroundTasks):
    body = await request.json()
    args = ["--override", _safe_arg(body["override"])] if body.get("override") else []
    background_tasks.add_task(run_agent, "rag_agent", "rag_agent.py", args)
    return {"status": "started", "agent": "rag_agent"}


@app.post("/api/agents/briefing")
async def trigger_briefing(request: Request, background_tasks: BackgroundTasks):
    body = await request.json()
    override = _safe_arg(body.get("override", ""))
    args = ["--override", override] if override else []
    background_tasks.add_task(run_agent, "briefing_agent", "briefing_agent.py", args, override)
    return {"status": "started", "agent": "briefing_agent"}


@app.post("/api/agents/content")
async def trigger_content(request: Request, background_tasks: BackgroundTasks):
    body = await request.json()
    override = _safe_arg(body.get("override", ""))
    args = ["--override", override] if override else []
    background_tasks.add_task(run_agent, "content_agent", "content_agent.py", args, override)
    return {"status": "started", "agent": "content_agent"}


class BriefRequest(BaseModel):
    client: str

@app.post("/api/agents/brief")
async def trigger_brief(req: BriefRequest, background_tasks: BackgroundTasks):
    background_tasks.add_task(run_agent, "coaching_brief", "coaching_brief.py", [req.client])
    return {"status": "started", "agent": "coaching_brief", "client": req.client}


# ── Deck Agent (Phase 2A) ─────────────────────────────────────────────────────

_DECK_TYPES = Literal["strategy", "pitch", "regulatory", "coaching", "ma", "briefing"]

class DeckRequest(BaseModel):
    topic: str
    deck_type: _DECK_TYPES = "strategy"
    client_name: str = ""
    context: str = ""

@app.post("/api/decks/generate")
async def generate_deck(req: DeckRequest, background_tasks: BackgroundTasks):
    args = [
        "--topic",   _safe_arg(req.topic),
        "--type",    _safe_arg(req.deck_type),
        "--context", _safe_arg(req.context),
    ]
    if req.client_name:
        args += ["--client", _safe_arg(req.client_name)]
    background_tasks.add_task(run_agent, "deck_agent", "deck_agent.py", args, req.topic)
    return {"status": "started", "agent": "deck_agent",
            "topic": req.topic, "deck_type": req.deck_type}

@app.get("/api/decks")
def list_decks():
    decks_dir = ATHENA / "documents" / "decks"
    if not decks_dir.exists():
        return {"decks": []}
    files = sorted(
        [f for f in decks_dir.glob("*.pptx")],
        key=lambda f: f.stat().st_mtime, reverse=True,
    )
    return {"decks": [{"filename": f.name, "path": str(f),
                       "modified": datetime.fromtimestamp(f.stat().st_mtime).isoformat()}
                      for f in files[:20]]}

@app.get("/api/decks/download/{filename}")
def download_deck(filename: str):
    decks_dir = (ATHENA / "documents" / "decks").resolve()
    path = (decks_dir / filename).resolve()
    if not str(path).startswith(str(decks_dir)) or not path.exists() or path.suffix != ".pptx":
        raise HTTPException(status_code=404, detail="Deck not found")
    return FileResponse(str(path), media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        filename=filename)

_SAFE_PPTX = re.compile(r'^[A-Za-z0-9_\-]+\.pptx$')

@app.get("/api/decks/{filename}/slides")
def preview_deck_slides(filename: str):
    """Extract slide text for in-browser preview — no external conversion needed."""
    if not _SAFE_PPTX.match(filename):
        raise HTTPException(404, "Not found")
    decks_dir = (ATHENA / "documents" / "decks").resolve()
    path = (decks_dir / filename).resolve()
    if not str(path).startswith(str(decks_dir)) or not path.exists():
        raise HTTPException(404, "Deck not found")
    from pptx import Presentation as _Prs
    prs = _Prs(str(path))
    slides = []
    for i, slide in enumerate(prs.slides):
        texts, title = [], ""
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                t = shape.text_frame.text.strip()
                if t:
                    if not title:
                        title = t
                    else:
                        texts.append(t)
            elif hasattr(shape, "table"):
                for row in shape.table.rows:
                    row_text = " | ".join(
                        cell.text_frame.text.strip() for cell in row.cells
                        if cell.text_frame.text.strip()
                    )
                    if row_text:
                        texts.append(row_text)
        slides.append({"index": i, "title": title, "texts": texts, "slide_number": i + 1})
    return {"slides": slides, "filename": filename, "total": len(slides)}

@app.post("/api/decks/bulk-delete")
async def bulk_delete_decks(request: Request):
    body = await request.json()
    filenames = body.get("filenames", [])
    decks_dir = (ATHENA / "documents" / "decks").resolve()
    deleted = []
    for fn in filenames:
        if not _SAFE_PPTX.match(fn):
            continue
        path = (decks_dir / fn).resolve()
        if str(path).startswith(str(decks_dir)) and path.exists():
            path.unlink()
            deleted.append(fn)
    _audit("decks_bulk_delete", f"count={len(deleted)}", request)
    return {"deleted": deleted}

@app.post("/api/decks/bulk-accept")
async def bulk_accept_decks(request: Request):
    body = await request.json()
    filenames = body.get("filenames", [])
    decks_dir = (ATHENA / "documents" / "decks").resolve()
    approved_dir = decks_dir / "approved"
    approved_dir.mkdir(exist_ok=True)
    accepted = []
    for fn in filenames:
        if not _SAFE_PPTX.match(fn):
            continue
        src = (decks_dir / fn).resolve()
        if str(src).startswith(str(decks_dir)) and src.exists():
            src.rename(approved_dir / fn)
            accepted.append(fn)
    _audit("decks_bulk_accept", f"count={len(accepted)}", request)
    return {"accepted": accepted}

@app.post("/api/decks/bulk-reject")
async def bulk_reject_decks(request: Request):
    body = await request.json()
    filenames = body.get("filenames", [])
    decks_dir = (ATHENA / "documents" / "decks").resolve()
    rejected = []
    for fn in filenames:
        if not _SAFE_PPTX.match(fn):
            continue
        path = (decks_dir / fn).resolve()
        if str(path).startswith(str(decks_dir)) and path.exists():
            path.unlink()
            rejected.append(fn)
    _audit("decks_bulk_reject", f"count={len(rejected)}", request)
    return {"rejected": rejected}

@app.post("/api/iso/lesson-deck/{lesson_filename}")
async def export_iso_lesson_deck(lesson_filename: str, background_tasks: BackgroundTasks):
    """Convert an ISO coaching lesson markdown into a management-consulting PPTX."""
    if not re.match(r'^[A-Za-z0-9_\-]+\.md$', lesson_filename):
        raise HTTPException(400, "Invalid filename")
    lesson_path: Optional[Path] = None
    for sub in ("iso14971", "iso13485"):
        candidate = ATHENA / "coaching" / sub / lesson_filename
        if candidate.exists():
            lesson_path = candidate
            break
    if not lesson_path:
        raise HTTPException(404, "Lesson not found")
    stem = lesson_path.stem
    is_14971 = stem.startswith("14971_")
    standard = "ISO 14971:2019" if is_14971 else "ISO 13485:2016"
    parts = stem.replace("14971_", "").replace("clause_", "").split("_")
    clause_parts, title_parts = [], []
    for p in parts:
        if re.match(r'^\d+$', p) and not title_parts:
            clause_parts.append(p)
        else:
            title_parts.append(p)
    clause = ".".join(clause_parts) if clause_parts else "?"
    title = " ".join(w.capitalize() for w in title_parts) if title_parts else stem
    background_tasks.add_task(_build_iso_pptx, lesson_path, standard, clause, title)
    return {"status": "started", "standard": standard, "clause": clause, "title": title}

def _build_iso_pptx(lesson_path: Path, standard: str, clause: str, title: str):
    """Parse ISO lesson markdown and build a branded PPTX. Runs in background."""
    try:
        import sys as _sys
        _sys.path.insert(0, str(AGENTS_DIR))
        from pptx import Presentation as _Prs
        from deck_skills import DeckSkills

        raw = lesson_path.read_text(encoding="utf-8", errors="replace")

        def _section(text: str, heading: str) -> str:
            pattern = rf"(?:^|\n)#+\s+{re.escape(heading)}[^\n]*\n(.*?)(?=\n#+\s|\Z)"
            m = re.search(pattern, text, re.IGNORECASE | re.DOTALL)
            return m.group(1).strip() if m else ""

        def _bullets_from(text: str) -> list:
            lines = []
            for ln in text.splitlines():
                ln = ln.strip()
                if ln.startswith(("- ", "* ", "• ")):
                    lines.append(ln[2:].strip())
                elif re.match(r'^\d+\.\s', ln):
                    lines.append(re.sub(r'^\d+\.\s*', '', ln))
            return lines[:6]

        overview      = _section(raw, "What This Clause") or _section(raw, "Overview") or raw[:400]
        req_text      = _section(raw, "Requirements") or _section(raw, "Key Requirements")
        reqs          = _bullets_from(req_text) or [
            f"Documented procedure for {title}",
            "Evidence of implementation and review",
            "Integration with QMS processes",
        ]
        mistakes_text = _section(raw, "Common Mistakes") or _section(raw, "Common Errors")
        mistakes      = _bullets_from(mistakes_text)
        career_text   = _section(raw, "Career Connection") or ""

        cs_blocks = re.findall(
            r'(?:^|\n)#+\s+(?:Case Study|Example)[^\n]*\n(.*?)(?=\n#+\s|\Z)',
            raw, re.IGNORECASE | re.DOTALL
        )

        terms: dict = {}
        for row in re.findall(r'\|([^|]+)\|([^|]+)\|', raw)[:8]:
            k, v = row[0].strip(), row[1].strip()
            if k and v and k.lower() not in ("term", "definition", "concept"):
                terms[k] = v
        if not terms:
            for k, v in re.findall(r'\*\*([^*]+)\*\*\s*[:\-—]\s*([^\n]{10,})', raw)[:6]:
                terms[k.strip()] = v.strip()

        prs = _Prs()
        prs.slide_width  = int(13.333 * 914400)
        prs.slide_height = int(7.5   * 914400)
        ds = DeckSkills(prs)

        ds.add_iso_cover(standard, clause, title, datetime.now().strftime("%B %Y"))
        ds.add_iso_overview(
            what_it_covers=overview[:500],
            why_it_matters=(
                f"Clause {clause} is a core requirement for {standard} certification. "
                "Non-conformance is a common audit finding."
            ),
        )
        ds.add_iso_requirements(reqs)

        for i, block in enumerate(cs_blocks[:2]):
            lines = [ln.strip() for ln in block.splitlines() if ln.strip()]
            company  = lines[0] if lines else f"Case Study {i+1}"
            scenario = lines[1] if len(lines) > 1 else block[:200]
            finding  = lines[2] if len(lines) > 2 else "Finding details in the full lesson."
            lesson   = lines[3] if len(lines) > 3 else "Apply rigorous documentation practices."
            is_real  = "real" in block[:120].lower() or "fda" in block[:60].lower()
            ds.add_iso_case_study_slide(
                title=f"Case Study {i+1}", company=company,
                scenario=scenario, finding=finding, lesson=lesson, is_real=is_real,
            )

        if mistakes:
            ds.add_iso_mistakes(mistakes)
        if terms or career_text:
            ds.add_iso_key_terms(terms, career_note=career_text[:400])

        decks_dir = ATHENA / "documents" / "decks"
        decks_dir.mkdir(parents=True, exist_ok=True)
        slug = re.sub(r'[^a-z0-9]+', '_', f"iso_{clause}_{title}".lower()).strip('_')[:40]
        out  = decks_dir / f"{datetime.now().strftime('%Y%m%d_%H%M%S')}_{slug}.pptx"
        prs.save(str(out))
    except Exception as exc:
        import logging as _log
        _log.getLogger("server").error(f"ISO deck build failed: {exc}", exc_info=True)


# ── LangGraph coaching orchestration (Phase 1A) ──────────────────────────────
async def _run_coaching_workflow(client: str):
    """Run the coaching graph to the human gate, streaming progress to the UI."""
    await manager.broadcast({"type": "agent_start", "agent": "orchestrator",
                             "ts": datetime.now().isoformat()})
    try:
        result = await asyncio.to_thread(_orchestrator().start_coaching, client)
        for step in result.get("steps", []):
            await manager.broadcast({"type": "agent_log", "agent": "orchestrator", "line": step})
        await manager.broadcast({"type": "agent_log", "agent": "orchestrator",
                                 "line": f"[HUMAN GATE] Review #{result.get('review_id')} "
                                         f"awaiting Steven's approval (thread {result['thread_id']})"})
        await manager.broadcast({"type": "agent_done", "agent": "orchestrator",
                                 "status": "awaiting_review", "ts": datetime.now().isoformat()})
    except Exception as e:
        await manager.broadcast({"type": "agent_log", "agent": "orchestrator", "line": f"[ERROR] {e}"})
        await manager.broadcast({"type": "agent_done", "agent": "orchestrator",
                                 "status": "error", "ts": datetime.now().isoformat()})

@app.post("/api/orchestrate/coaching")
async def orchestrate_coaching(req: BriefRequest, background_tasks: BackgroundTasks):
    """Run the full coaching workflow (brief -> human gate -> finalize) via LangGraph."""
    background_tasks.add_task(_run_coaching_workflow, req.client)
    return {"status": "started", "workflow": "coaching", "client": req.client}


class ISORequest(BaseModel):
    clause: Optional[str] = None

@app.post("/api/agents/iso")
async def trigger_iso(req: ISORequest, background_tasks: BackgroundTasks, request: Request):
    clause_detail = req.clause or "next"
    _audit("agent_trigger", f"agent=iso_coach clause={clause_detail}", request)
    args = ["--clause", req.clause] if req.clause else ["--next"]
    background_tasks.add_task(run_agent, "iso_coach", "iso_coach_agent.py", args)
    return {"status": "started", "agent": "iso_coach"}


@app.post("/api/agents/voice-optimizer")
async def trigger_voice_optimizer(request: Request, background_tasks: BackgroundTasks):
    """
    Run one PDCA optimization cycle for Athena's voice pipeline.
    Tunes wake_threshold, silence params, Whisper model, and voice_assistant
    system prompt — all within safe bounds. Code-level improvements are queued
    to the review queue for Steven's approval.
    """
    try:
        body = await request.json()
    except Exception:
        body = {}
    dry_run  = body.get("dry_run", False)
    args     = ["--dry-run"] if dry_run else []
    background_tasks.add_task(run_agent, "voice_optimizer", "voice_optimizer.py", args)
    return {"status": "started", "agent": "voice_optimizer", "dry_run": dry_run}


# ── Document generation ───────────────────────────────────────────────────────

class DocRequest(BaseModel):
    title:    str
    content:  str
    doc_type: str = "article"  # article | brief | report | sop


@app.post("/api/documents/generate")
async def generate_document(req: DocRequest):
    """Generate a branded .docx from content."""
    try:
        out_path = await asyncio.get_event_loop().run_in_executor(
            None, generate_docx, req.title, req.content, req.doc_type
        )
        return {"status": "ok", "path": str(out_path), "filename": out_path.name}
    except Exception as e:
        return JSONResponse(status_code=500, content={"error": str(e)})


# ── Figure embedding (Markdown ![alt](path) inside document content) ──────────
FIGURES_DIR      = ATHENA / "documents" / "figures"
MAX_FIG_WIDTH_IN = 6.1   # Letter (8.5in) minus 1.2in margins on each side


def _resolve_figure(src: str) -> Optional[Path]:
    """Resolve a figure referenced by ![alt](src). Tries absolute, the figures
    dir, Athena-relative, then cwd. Returns None if nothing exists."""
    cand = Path(src)
    candidates = [cand] if cand.is_absolute() else [
        FIGURES_DIR / src, ATHENA / "documents" / src, ATHENA / src, Path.cwd() / src,
    ]
    for c in candidates:
        try:
            if c.is_file():
                return c
        except OSError:
            continue
    return None


def generate_docx(title: str, content: str, doc_type: str) -> Path:
    """Generate branded .docx using python-docx.

    Content is Markdown. Supported block constructs:
      ``# / ## / ###``     headings
      ``- `` / ``* ``      bullets
      ``**text**``         bold (line or inline)
      ``---``              horizontal rule
      ``![caption](path)`` figure (chart/diagram), centered, width-constrained,
                           caption beneath. Render charts with
                           ``figures.bar_chart(...)`` and reference the path.
      ``> text``           callout / pull-stat box (shaded, left accent bar);
                           consecutive ``>`` lines merge into one box.
      ``| a | b |``        pipe table; a ``|---|`` separator row marks a header.
    """
    try:
        from docx import Document
        from docx.shared import Pt, RGBColor, Inches, Cm
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.enum.style import WD_STYLE_TYPE
        from docx.enum.table import WD_TABLE_ALIGNMENT
        from docx.oxml.ns import qn
        from docx.oxml import OxmlElement
        import re
    except ImportError:
        raise ImportError("python-docx not installed. Run: pip install python-docx")

    # Latitude MedTech colors
    LM_BLACK = RGBColor(0x1A, 0x1A, 0x1A)
    LM_SLATE = RGBColor(0x2C, 0x3E, 0x50)
    LM_BLUE  = RGBColor(0x5B, 0x7F, 0xA6)
    LM_MUTED = RGBColor(0x8A, 0x86, 0x80)
    LM_WHITE = RGBColor(0xFF, 0xFF, 0xFF)

    doc = Document()

    # Page margins
    for section in doc.sections:
        section.top_margin    = Inches(1.0)
        section.bottom_margin = Inches(1.0)
        section.left_margin   = Inches(1.2)
        section.right_margin  = Inches(1.2)

    # Default font
    doc.styles['Normal'].font.name = 'Calibri'
    doc.styles['Normal'].font.size = Pt(10.5)
    doc.styles['Normal'].font.color.rgb = LM_SLATE

    # ── Header ────────────────────────────────────────────────────────────────
    header = doc.sections[0].header
    header.paragraphs[0].clear()
    run = header.paragraphs[0].add_run("LATITUDE MEDTECH LLC")
    run.font.name  = 'Calibri'
    run.font.size  = Pt(7.5)
    run.font.color.rgb = LM_BLUE
    run.font.bold  = True
    header.paragraphs[0].paragraph_format.alignment = WD_ALIGN_PARAGRAPH.RIGHT

    # ── Footer ────────────────────────────────────────────────────────────────
    footer = doc.sections[0].footer
    footer.paragraphs[0].clear()
    run = footer.paragraphs[0].add_run(
        f"Latitude MedTech LLC  ·  San Diego, CA  ·  latitudemedtech.com  ·  "
        f"Generated {datetime.now().strftime('%B %d, %Y')}  ·  "
        "Educational purposes only — not regulatory advice"
    )
    run.font.name  = 'Calibri'
    run.font.size  = Pt(7)
    run.font.color.rgb = LM_MUTED
    footer.paragraphs[0].paragraph_format.alignment = WD_ALIGN_PARAGRAPH.CENTER

    # ── Title block ───────────────────────────────────────────────────────────
    title_para = doc.add_paragraph()
    title_para.paragraph_format.space_before = Pt(0)
    title_para.paragraph_format.space_after  = Pt(6)
    title_run = title_para.add_run(title.upper())
    title_run.font.name  = 'Calibri'
    title_run.font.size  = Pt(18)
    title_run.font.bold  = True
    title_run.font.color.rgb = LM_BLACK
    title_para.paragraph_format.alignment = WD_ALIGN_PARAGRAPH.LEFT

    # Accent line under title
    meta_para = doc.add_paragraph()
    meta_run  = meta_para.add_run(
        f"Latitude MedTech LLC  ·  {doc_type.title()}  ·  "
        f"{datetime.now().strftime('%B %d, %Y')}"
    )
    meta_run.font.name  = 'Calibri'
    meta_run.font.size  = Pt(8.5)
    meta_run.font.color.rgb = LM_BLUE
    meta_para.paragraph_format.space_after = Pt(18)

    # ── Content ───────────────────────────────────────────────────────────────
    def emit_runs(p, text):
        """Add runs to paragraph p, honoring inline **bold**."""
        for j, part in enumerate(re.split(r'\*\*(.+?)\*\*', text)):
            if not part:
                continue
            r = p.add_run(part)
            r.font.name = 'Calibri'; r.font.size = Pt(10.5)
            if j % 2 == 1:
                r.font.bold = True

    def emit_callout(text_lines):
        """Shaded pull-stat / callout box with a left accent bar (from > lines)."""
        p = doc.add_paragraph()
        pPr = p._p.get_or_add_pPr()
        pBdr = OxmlElement('w:pBdr')
        for edge, sz, col in (('left', '20', '5B7FA6'), ('top', '4', 'D9E1EA'),
                              ('bottom', '4', 'D9E1EA'), ('right', '4', 'D9E1EA')):
            e = OxmlElement('w:' + edge)
            e.set(qn('w:val'), 'single'); e.set(qn('w:sz'), sz)
            e.set(qn('w:space'), '8');    e.set(qn('w:color'), col)
            pBdr.append(e)
        pPr.append(pBdr)
        shd = OxmlElement('w:shd')
        shd.set(qn('w:val'), 'clear'); shd.set(qn('w:fill'), 'EFF3F7')
        pPr.append(shd)
        pf = p.paragraph_format
        pf.left_indent = Inches(0.14); pf.right_indent = Inches(0.10)
        pf.space_before = Pt(8);       pf.space_after = Pt(10)
        for k, t in enumerate(text_lines):
            if k:
                p.add_run().add_break()
            emit_runs(p, t)

    def shade_cell(cell, fill):
        tcPr = cell._tc.get_or_add_tcPr()
        shd = OxmlElement('w:shd')
        shd.set(qn('w:val'), 'clear'); shd.set(qn('w:fill'), fill)
        tcPr.append(shd)

    def set_table_borders(table):
        tblPr = table._tbl.tblPr
        borders = OxmlElement('w:tblBorders')
        for edge in ('top', 'left', 'bottom', 'right', 'insideH', 'insideV'):
            e = OxmlElement('w:' + edge)
            e.set(qn('w:val'), 'single'); e.set(qn('w:sz'), '4')
            e.set(qn('w:space'), '0');    e.set(qn('w:color'), 'CCCCCC')
            borders.append(e)
        ref = tblPr.find(qn('w:tblLook'))
        if ref is None:
            ref = tblPr.find(qn('w:tblCellMar'))
        if ref is not None:
            ref.addprevious(borders)
        else:
            tblPr.append(borders)

    def emit_table(block):
        """Render a Markdown pipe table; a |---| separator row marks a header."""
        grid = [[c.strip() for c in row.strip().strip('|').split('|')] for row in block]
        is_sep = lambda cells: all(set(c) <= set('-: ') and '-' in c
                                   for c in cells if c != '')
        has_header = len(grid) >= 2 and is_sep(grid[1])
        body = ([grid[0]] + grid[2:]) if has_header else grid
        ncols = max(len(r) for r in body)
        table = doc.add_table(rows=0, cols=ncols)
        table.alignment = WD_TABLE_ALIGNMENT.CENTER
        table.allow_autofit = False
        set_table_borders(table)
        col_w = Inches(MAX_FIG_WIDTH_IN / ncols)
        for ri, row in enumerate(body):
            cells = table.add_row().cells
            for ci in range(ncols):
                cell = cells[ci]; cell.width = col_w
                para = cell.paragraphs[0]
                para.paragraph_format.space_after = Pt(2)
                text = row[ci] if ci < len(row) else ''
                if has_header and ri == 0:
                    shade_cell(cell, '5B7FA6')
                    r = para.add_run(re.sub(r'\*\*(.+?)\*\*', r'\1', text))
                    r.font.name = 'Calibri'; r.font.size = Pt(10)
                    r.font.bold = True;      r.font.color.rgb = LM_WHITE
                else:
                    if has_header and ri % 2 == 0:
                        shade_cell(cell, 'F2F5F8')
                    emit_runs(para, text)
        doc.add_paragraph().paragraph_format.space_after = Pt(4)

    lines = content.splitlines()
    n    = len(lines)
    idx  = 0
    while idx < n:
        line = lines[idx].strip()

        if not line:
            doc.add_paragraph()
            idx += 1
            continue

        # Pipe table — consume the contiguous block of | … | lines
        if line.startswith('|'):
            block = []
            while idx < n and lines[idx].strip().startswith('|'):
                block.append(lines[idx].strip())
                idx += 1
            emit_table(block)
            continue

        # Callout / pull-stat — consume contiguous > lines
        if line.startswith('>'):
            quote = []
            while idx < n and lines[idx].strip().startswith('>'):
                quote.append(lines[idx].strip()[1:].strip())
                idx += 1
            emit_callout(quote)
            continue

        # H1
        if line.startswith('# '):
            p    = doc.add_heading(line[2:], level=1)
            p.runs[0].font.color.rgb = LM_BLACK
            p.runs[0].font.name = 'Calibri'
            idx += 1; continue

        # H2
        if line.startswith('## '):
            p    = doc.add_heading(line[3:], level=2)
            p.runs[0].font.color.rgb = LM_BLUE
            p.runs[0].font.name = 'Calibri'
            idx += 1; continue

        # H3
        if line.startswith('### '):
            p    = doc.add_heading(line[4:], level=3)
            p.runs[0].font.color.rgb = LM_SLATE
            p.runs[0].font.name = 'Calibri'
            idx += 1; continue

        # Bullet (inline bold supported)
        if line.startswith('- ') or line.startswith('* '):
            p = doc.add_paragraph(style='List Bullet')
            emit_runs(p, line[2:])
            idx += 1; continue

        # Whole-line bold paragraph (single bold span only)
        if line.startswith('**') and line.endswith('**') and line.count('**') == 2:
            p   = doc.add_paragraph()
            run = p.add_run(line[2:-2])
            run.font.bold = True
            run.font.name = 'Calibri'
            run.font.size = Pt(10.5)
            idx += 1; continue

        # Horizontal rule
        if line.startswith('---'):
            p = doc.add_paragraph()
            pPr = p._p.get_or_add_pPr()
            pBdr = OxmlElement('w:pBdr')
            bottom = OxmlElement('w:bottom')
            bottom.set(qn('w:val'), 'single')
            bottom.set(qn('w:sz'), '4')
            bottom.set(qn('w:space'), '1')
            bottom.set(qn('w:color'), '5B7FA6')
            pBdr.append(bottom)
            pPr.append(pBdr)
            idx += 1; continue

        # Figure — ![caption](path) on its own line
        img = re.match(r'^!\[(.*?)\]\((.+?)\)\s*$', line)
        if img:
            alt, src = img.group(1).strip(), img.group(2).strip()
            fig_path = _resolve_figure(src)
            if fig_path is None:
                miss = doc.add_paragraph()
                mrun = miss.add_run(f"[missing figure: {src}]")
                mrun.font.name = 'Calibri'; mrun.font.size = Pt(9)
                mrun.font.italic = True;    mrun.font.color.rgb = LM_MUTED
                idx += 1; continue
            # Constrain to content width; never upscale a smaller image
            width_in = MAX_FIG_WIDTH_IN
            try:
                from PIL import Image as _PILImage
                with _PILImage.open(fig_path) as im:
                    width_in = min(MAX_FIG_WIDTH_IN, im.width / 96.0)
            except Exception:
                pass
            doc.add_picture(str(fig_path), width=Inches(width_in))
            pic = doc.paragraphs[-1]
            pic.alignment = WD_ALIGN_PARAGRAPH.CENTER
            pic.paragraph_format.space_before = Pt(6)
            pic.paragraph_format.space_after  = Pt(2)
            if alt:
                cap = doc.add_paragraph()
                cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
                cap.paragraph_format.space_after = Pt(10)
                crun = cap.add_run(alt)
                crun.font.name = 'Calibri'; crun.font.size = Pt(8.5)
                crun.font.italic = True;    crun.font.color.rgb = LM_MUTED
            idx += 1; continue

        # Normal paragraph — inline bold
        p = doc.add_paragraph()
        p.paragraph_format.space_after = Pt(4)
        emit_runs(p, line)
        idx += 1

    # ── Disclaimer ────────────────────────────────────────────────────────────
    doc.add_paragraph()
    disc = doc.add_paragraph()
    disc.paragraph_format.space_before = Pt(12)
    disc_run = disc.add_run(
        "DISCLAIMER: This document is produced by Latitude MedTech LLC for educational and planning purposes only. "
        "Nothing herein constitutes formal regulatory, legal, or compliance advice. "
        "Always defer to official regulations and guidance from the FDA, ISO, EU Commission, "
        "and applicable regulatory bodies."
    )
    disc_run.font.name   = 'Calibri'
    disc_run.font.size   = Pt(8)
    disc_run.font.italic = True
    disc_run.font.color.rgb = LM_MUTED

    # ── Save ──────────────────────────────────────────────────────────────────
    out_dir  = ATHENA / 'documents'
    out_dir.mkdir(exist_ok=True)
    slug     = re.sub(r'[^a-z0-9]+', '_', title.lower())[:50]
    out_path = out_dir / f"{datetime.now().strftime('%Y-%m-%d')}_{slug}.docx"
    doc.save(str(out_path))
    return out_path


@app.get("/api/documents")
def list_documents():
    # All agent output folders surface here so every deliverable is visible
    # in the Documents hub — not stranded on disk.
    sources = [
        ("documents",       ATHENA / "documents",                    "*.docx", False),
        ("sow",             ATHENA / "documents" / "sow",            "*.docx", False),
        ("regulatory",      ATHENA / "documents" / "regulatory",     "*.docx", False),
        ("ma_intelligence", ATHENA / "ops" / "ma_intelligence",      "*.md",   True),
        ("consulting",      ATHENA / "ops" / "consulting",           "*.md",   True),
        ("marketing",       ATHENA / "ops" / "marketing",            "*.md",   True),
        ("briefings",       ATHENA / "briefings",                    "*.md",   True),
        ("content",         ATHENA / "content" / "drafts",           "*.md",   False),
        ("iso13485",        ATHENA / "coaching" / "iso13485",        "*.md",   False),
        ("iso14971",        ATHENA / "coaching" / "iso14971",        "*.md",   False),
        ("briefs",          ATHENA / "coaching" / "briefs",          "*.md",   False),
    ]
    items = []
    for folder, base, pattern, include_summary in sources:
        if not base.exists():
            continue
        for f in base.glob(pattern):
            _st = f.stat()
            item = {
                "filename": f.name,
                "folder": folder,
                "path": str(f),
                "modified": datetime.fromtimestamp(_st.st_mtime).isoformat(),
                "mtime": _st.st_mtime,
            }
            if include_summary:
                item["summary"] = _extract_md_summary(f)
            items.append(item)
    items.sort(key=lambda d: d["mtime"], reverse=True)
    for d in items:
        d.pop("mtime", None)
    return {"documents": items[:50]}


@app.get("/api/documents/open/{filename}")
def open_document(filename: str, folder: str = "documents"):
    base = FOLDER_MAP.get(folder, ATHENA / "documents")
    f = (base / _safe_filename(filename)).resolve()
    if not str(f).startswith(str(base.resolve())) or not f.exists():
        return JSONResponse(status_code=404, content={"error": "Not found"})
    os.startfile(str(f))
    return {"status": "opened", "filename": filename}


# ── WebSocket ─────────────────────────────────────────────────────────────────

@app.websocket("/ws")
async def websocket_endpoint(ws: WebSocket, token: str = ""):
    # SOC II CC6.6 — reject unauthenticated WebSocket upgrades
    if not secrets.compare_digest(token.encode(), SESSION_TOKEN.encode()):
        _audit("ws_auth_failure", "rejected /ws without valid token")
        await ws.close(code=4001)
        return
    global _goodbye_task
    # Cancel any pending goodbye — this is a reconnect (e.g. page refresh).
    if _goodbye_task and not _goodbye_task.done():
        _goodbye_task.cancel()
    await manager.connect(ws)
    try:
        while True:
            await ws.receive_text()
    except WebSocketDisconnect:
        manager.disconnect(ws)
        if not manager.active and not _shutting_down:
            _goodbye_task = asyncio.get_event_loop().create_task(_ws_goodbye())


@app.websocket("/ws/voice")
async def voice_ws_endpoint(ws: WebSocket, token: str = ""):
    # SOC II CC6.6 — reject unauthenticated voice WebSocket upgrades
    if not secrets.compare_digest(token.encode(), SESSION_TOKEN.encode()):
        _audit("ws_auth_failure", "rejected /ws/voice without valid token")
        await ws.close(code=4001)
        return
    if VOICE_AVAILABLE:
        await voice_websocket_endpoint(ws)
    else:
        await ws.accept()
        await ws.send_json({"type": "voice_error", "message": "Voice bridge not available"})
        await ws.close()




# ── Settings routes ───────────────────────────────────────────────────────────

sys.path.insert(0, str(AGENTS_DIR))
try:
    from settings_manager import settings as agent_settings
    SETTINGS_AVAILABLE = True
except ImportError:
    SETTINGS_AVAILABLE = False


@app.get("/api/settings")
def get_settings():
    if not SETTINGS_AVAILABLE:
        return JSONResponse(status_code=503, content={"error": "settings_manager.py not found in agents folder"})
    return agent_settings.get_all()


@app.post("/api/settings")
async def update_settings(request: Request):
    if not SETTINGS_AVAILABLE:
        return JSONResponse(status_code=503, content={"error": "settings_manager not available"})
    try:
        body = await request.json()
        for key_path, value in body.items():
            agent_settings.set(key_path, value)
        _audit("settings_update", f"keys={list(body.keys())}", request)
        return {"status": "saved", "updated_at": datetime.now().isoformat()}
    except Exception as e:
        return JSONResponse(status_code=500, content={"error": str(e)})


@app.post("/api/settings/prompt/{agent_name}")
async def update_prompt(agent_name: str, request: Request):
    if not SETTINGS_AVAILABLE:
        return JSONResponse(status_code=503, content={"error": "settings_manager not available"})
    try:
        body  = await request.json()
        prompt = body.get("prompt", "")
        agent_settings.set_prompt(agent_name, prompt)
        _audit("settings_prompt_update", f"agent={agent_name}", request)
        return {"status": "saved", "agent": agent_name}
    except Exception as e:
        return JSONResponse(status_code=500, content={"error": str(e)})


@app.post("/api/settings/reset")
def reset_settings():
    if not SETTINGS_AVAILABLE:
        return JSONResponse(status_code=503, content={"error": "settings_manager not available"})
    agent_settings.reset_to_defaults()
    _audit("settings_reset", "all settings restored to defaults")
    return {"status": "reset"}


# ── File management routes ────────────────────────────────────────────────────

_HOME_ATHENA = ATHENA   # canonical code-tree root; agents write here
FOLDER_MAP = {
    "briefings":       _HOME_ATHENA / "briefings",
    "content":         _HOME_ATHENA / "content" / "drafts",
    "content/drafts":  _HOME_ATHENA / "content" / "drafts",   # legacy key
    "briefs":          _HOME_ATHENA / "coaching" / "briefs",
    "coaching/briefs": _HOME_ATHENA / "coaching" / "briefs",  # legacy key
    "documents":       _HOME_ATHENA / "documents",
    "sow":             _HOME_ATHENA / "documents" / "sow",
    "regulatory":      _HOME_ATHENA / "documents" / "regulatory",
    "iso13485":        _HOME_ATHENA / "coaching" / "iso13485",
    "iso14971":        _HOME_ATHENA / "coaching" / "iso14971",
    "marketing":       _HOME_ATHENA / "ops" / "marketing",
    "learning":        _HOME_ATHENA / "knowledge_base" / "learning",
    "hr":              _HOME_ATHENA / "ops" / "hr",
    "consulting":      _HOME_ATHENA / "ops" / "consulting",
    "ma_intelligence": _HOME_ATHENA / "ops" / "ma_intelligence",
}


@app.post("/api/files/delete")
async def delete_file(request: Request):
    try:
        body     = await request.json()
        folder   = body.get("folder", "")
        filename = body.get("filename", "")
        base     = FOLDER_MAP.get(folder)
        if not base:
            print(f"DELETE ERROR: Unknown folder '{folder}'. Known: {list(FOLDER_MAP.keys())}")
            return JSONResponse(status_code=400, content={"error": f"Unknown folder: {folder}"})
        f = base / filename
        print(f"DELETE: {f} exists={f.exists()}")
        if f.exists():
            f.unlink()
            _audit("file_delete", f"filename={filename} folder={folder}", request)
            return {"status": "deleted", "filename": filename}
        # Try searching all folders
        for fname, fbase in FOLDER_MAP.items():
            candidate = fbase / filename
            if candidate.exists():
                candidate.unlink()
                _audit("file_delete", f"filename={filename} folder={fname}", request)
                return {"status": "deleted", "filename": filename, "found_in": fname}
        return JSONResponse(status_code=404, content={"error": f"File not found: {filename} in {base}"})
    except Exception as e:
        print(f"DELETE EXCEPTION: {e}")
        return JSONResponse(status_code=500, content={"error": str(e)})


@app.post("/api/files/delete-bulk")
async def delete_files_bulk(request: Request):
    try:
        body  = await request.json()
        items = body.get("items", [])   # [{folder, filename}, ...]
        deleted, failed = [], []
        for item in items:
            folder   = item.get("folder", "")
            filename = item.get("filename", "")
            base     = FOLDER_MAP.get(folder)
            deleted_this = False
            if base:
                f = base / filename
                if f.exists():
                    f.unlink()
                    deleted.append(filename)
                    deleted_this = True
            if not deleted_this:
                # Search all folders
                for fbase in FOLDER_MAP.values():
                    candidate = fbase / filename
                    if candidate.exists():
                        candidate.unlink()
                        deleted.append(filename)
                        deleted_this = True
                        break
            if not deleted_this:
                failed.append(filename)
        _audit("file_delete_bulk", f"deleted={len(deleted)} failed={len(failed)}", request)
        return {"deleted": deleted, "failed": failed, "count": len(deleted)}
    except Exception as e:
        return JSONResponse(status_code=500, content={"error": str(e)})


@app.post("/api/files/save")
async def save_file(request: Request):
    try:
        body     = await request.json()
        filename = body.get("filename", "")
        content  = body.get("content", "")
        for folder, base in FOLDER_MAP.items():
            f = base / filename
            if f.exists():
                f.write_text(content, encoding="utf-8")
                _audit("file_save", f"filename={filename} folder={folder}", request)
                return {"status": "saved", "filename": filename}
        return JSONResponse(status_code=404, content={"error": f"File not found in any folder: {filename}"})
    except Exception as e:
        print(f"SAVE EXCEPTION: {e}")
        return JSONResponse(status_code=500, content={"error": str(e)})


@app.get("/api/coaching/briefs")
def list_briefs():
    briefs_dir = _HOME_ATHENA / "coaching" / "briefs"
    if not briefs_dir.exists():
        return {"briefs": []}
    files = sorted(briefs_dir.glob("*.md"), key=lambda f: f.stat().st_mtime, reverse=True)
    return {"briefs": [
        {"filename": f.name, "modified": datetime.fromtimestamp(f.stat().st_mtime).isoformat()}
        for f in files[:20]
    ]}


@app.get("/api/coaching/briefs/{filename}")
def get_brief(filename: str):
    _safe_filename(filename)   # SOC II CC6.1 — prevent path traversal
    briefs_dir = _HOME_ATHENA / "coaching" / "briefs"
    f = (briefs_dir / filename).resolve()
    if not str(f).startswith(str(briefs_dir.resolve())) or not f.exists():
        return JSONResponse(status_code=404, content={"error": "Not found"})
    return {"filename": filename, "content": f.read_text(encoding="utf-8", errors="replace")}


# ── ISO 13485 coach content ───────────────────────────────────────────────────

@app.get("/api/iso/lessons")
def list_iso_lessons():
    all_files = []
    for sub in ("iso13485", "iso14971"):
        d = _HOME_ATHENA / "coaching" / sub
        if d.exists():
            all_files.extend(d.glob("*.md"))
    if not all_files:
        return {"lessons": []}
    files = sorted(all_files, key=lambda f: f.stat().st_mtime, reverse=True)

    def _parse_iso_stem(stem: str):
        # ISO 14971 filenames start with "14971_": e.g. "14971_5_2_risk_analysis"
        if stem.startswith("14971_"):
            rest  = stem[len("14971_"):]
            parts = rest.split("_")
            num_parts, word_parts = [], []
            for p in parts:
                if p.replace(".", "").isdigit():
                    num_parts.append(p)
                else:
                    word_parts.append(p)
            clause = "14971." + ".".join(num_parts) if num_parts else "14971." + rest
            title  = " ".join(w.capitalize() for w in word_parts) if word_parts else clause
            return clause, title
        # ISO 13485: "clause_4_1_general_qms_requirements"
        parts = stem.split("_")
        start = 1 if parts and parts[0].lower() == "clause" else 0
        parts = parts[start:]
        num_parts, word_parts = [], []
        for p in parts:
            if p.replace(".", "").isdigit():
                num_parts.append(p)
            else:
                word_parts.append(p)
        clause = ".".join(num_parts) if num_parts else stem
        title  = " ".join(w.capitalize() for w in word_parts) if word_parts else clause
        return clause, title

    lessons = []
    for f in files[:100]:
        clause, title = _parse_iso_stem(f.stem)
        lessons.append({
            "filename": f.name,
            "clause":   clause,
            "title":    title,
            "standard": "ISO 14971" if f.name.startswith("14971_") else "ISO 13485",
            "modified": datetime.fromtimestamp(f.stat().st_mtime).isoformat(),
        })
    return {"lessons": lessons}


@app.get("/api/iso/lessons/{filename}")
def get_iso_lesson(filename: str):
    _safe_filename(filename)   # SOC II CC6.1 — prevent path traversal
    for sub in ("iso14971", "iso13485"):
        base = (_HOME_ATHENA / "coaching" / sub).resolve()
        f    = (base / filename).resolve()
        if str(f).startswith(str(base)) and f.exists():
            return {"filename": filename, "content": f.read_text(encoding="utf-8", errors="replace")}
    return JSONResponse(status_code=404, content={"error": "Not found"})


# ── File serve (for in-browser viewer) ───────────────────────────────────────

from fastapi.responses import Response as FResponse

@app.get("/api/files/serve/{folder}/{filename}")
def serve_file(folder: str, filename: str):
    """Serve raw file bytes for in-browser viewing (PDF / DOCX / MD)."""
    _safe_filename(filename)                        # raises 400 on traversal
    base = FOLDER_MAP.get(folder)
    if not base:
        return JSONResponse(status_code=400, content={"error": "Unknown folder"})
    f = (base / filename).resolve()
    # Confirm resolved path is still inside the expected base directory
    if not str(f).startswith(str(base.resolve())):
        raise HTTPException(status_code=403, detail="Access denied")
    if not f.exists():
        return JSONResponse(status_code=404, content={"error": "Not found"})
    suffix = f.suffix.lower()
    mime = {
        ".pdf":  "application/pdf",
        ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ".md":   "text/plain; charset=utf-8",
        ".txt":  "text/plain; charset=utf-8",
    }.get(suffix, "application/octet-stream")
    return FResponse(content=f.read_bytes(), media_type=mime,
                     headers={"Content-Disposition": f'inline; filename="{filename}"'})


# ── Google Drive viewer ───────────────────────────────────────────────────────
#
# Config (set in Athena/voice/.env):
#   GOOGLE_CREDENTIALS_PATH  — absolute path to service account credentials.json
#                              (default: Athena/credentials.json)
#   GOOGLE_DRIVE_FOLDER_ID   — ID of the Drive folder shared with the service account
#   GOOGLE_SHARE_EMAIL       — Google account to share previews with
#
# Upload happens once per file; cache stored in Athena/drive_cache.json.
# Re-opening an unchanged file costs zero — cached Drive file ID returned instantly.

_DRIVE_CACHE_PATH   = ATHENA / "drive_cache.json"
_GDRIVE_SECRETS     = str(ATHENA / "client_secrets.json")
_GDRIVE_TOKEN       = str(ATHENA / "drive_token.json")
_GDRIVE_SCOPES      = ["https://www.googleapis.com/auth/drive"]
_GDRIVE_FOLDER_ID   = os.getenv("GOOGLE_DRIVE_FOLDER_ID", "")  # optional; uploads go to Drive root if blank

# Thread state for the non-blocking OAuth flow
_drive_auth_state: dict = {"running": False, "error": ""}
_drive_auth_lock  = threading.Lock()   # guards the check-then-set in google_auth_connect

# Local MIME → (upload MIME, target Google MIME)
_TO_GOOGLE_MIME = {
    ".docx": ("application/vnd.openxmlformats-officedocument.wordprocessingml.document",
              "application/vnd.google-apps.document"),
    ".pptx": ("application/vnd.openxmlformats-officedocument.presentationml.presentation",
              "application/vnd.google-apps.presentation"),
    ".pdf":  ("application/pdf", None),   # keep as PDF; Drive viewer handles it natively
}

# Embed URL per Google MIME
_EMBED_URL = {
    "application/vnd.google-apps.document":     "https://docs.google.com/document/d/{id}/preview",
    "application/vnd.google-apps.presentation": "https://docs.google.com/presentation/d/{id}/preview",
    None: "https://drive.google.com/file/d/{id}/preview",
}


def _drive_cache_load() -> dict:
    try:
        return json.loads(_DRIVE_CACHE_PATH.read_text(encoding="utf-8"))
    except Exception:
        return {}


def _drive_cache_save(cache: dict):
    try:
        _DRIVE_CACHE_PATH.write_text(json.dumps(cache, indent=2), encoding="utf-8")
    except Exception:
        pass


def _build_drive_service():
    """Return an authenticated Google Drive v3 service (OAuth 2.0), or None if not ready."""
    token_path   = Path(_GDRIVE_TOKEN)
    secrets_path = Path(_GDRIVE_SECRETS)
    if not secrets_path.exists():
        return None
    try:
        from google.oauth2.credentials import Credentials
        from google.auth.transport.requests import Request
        from googleapiclient.discovery import build as _gbuild
        creds = None
        if token_path.exists():
            creds = Credentials.from_authorized_user_file(str(token_path), _GDRIVE_SCOPES)
        if not creds or not creds.valid:
            if creds and creds.expired and creds.refresh_token:
                creds.refresh(Request())
                token_path.write_text(creds.to_json())
            else:
                return None  # needs browser auth — call POST /api/google/auth
        return _gbuild("drive", "v3", credentials=creds, cache_discovery=False)
    except Exception as exc:
        print(f"[drive] Could not build service: {exc}")
        return None


@app.post("/api/google/auth")
def google_auth_connect():
    """
    Spawn a background thread to run the OAuth 2.0 browser flow and save
    drive_token.json.  Returns {"started": true} immediately so the HTTP
    request doesn't hang while the user completes consent in the browser.
    Poll GET /api/google/auth-status every 2 s until configured=true or reason=error.
    """
    global _drive_auth_state
    secrets_path = Path(_GDRIVE_SECRETS)
    if not secrets_path.exists():
        raise HTTPException(503, "client_secrets.json not found in Athena root folder")
    with _drive_auth_lock:
        if _drive_auth_state.get("running"):
            return {"started": True, "message": "Authorization already in progress…"}
        _drive_auth_state = {"running": True, "error": ""}

    def _run():
        global _drive_auth_state
        try:
            from google_auth_oauthlib.flow import InstalledAppFlow
            flow  = InstalledAppFlow.from_client_secrets_file(str(secrets_path), _GDRIVE_SCOPES)
            creds = flow.run_local_server(port=0, open_browser=True)
            Path(_GDRIVE_TOKEN).write_text(creds.to_json())
            _drive_auth_state = {"running": False, "error": ""}
        except Exception as exc:
            _drive_auth_state = {"running": False, "error": str(exc)}
            print(f"[drive-auth] OAuth flow failed: {exc}")

    threading.Thread(target=_run, daemon=True).start()
    return {"started": True, "message": "Browser opening for Google authorization…"}


@app.get("/api/google/auth-status")
def google_auth_status():
    """Check whether Drive credentials exist and are valid (or auto-refreshable)."""
    global _drive_auth_state
    secrets_path = Path(_GDRIVE_SECRETS)
    token_path   = Path(_GDRIVE_TOKEN)

    # Expose thread state first
    if _drive_auth_state.get("running"):
        return {"configured": False, "reason": "in_progress",
                "message": "Authorization in progress — complete consent in your browser."}
    err = _drive_auth_state.pop("error", "")   # read + clear atomically; "" → not an error
    if err:
        return {"configured": False, "reason": "error", "message": err}

    if not secrets_path.exists():
        return {"configured": False, "reason": "client_secrets.json missing"}
    if not token_path.exists():
        return {"configured": False, "reason": "not_authorized",
                "message": "Authorization required — click 'Connect Google Drive'."}
    try:
        from google.oauth2.credentials import Credentials
        from google.auth.transport.requests import Request
        creds = Credentials.from_authorized_user_file(str(token_path), _GDRIVE_SCOPES)
        if creds.valid:
            return {"configured": True, "reason": "Token valid"}
        if creds.expired and creds.refresh_token:
            creds.refresh(Request())
            token_path.write_text(creds.to_json())
            return {"configured": True, "reason": "Token refreshed"}
        return {"configured": False, "reason": "Token expired — re-authorization required"}
    except Exception as exc:
        return {"configured": False, "reason": str(exc)}


@app.get("/api/files/google-view/{folder}/{filename}")
def google_view(folder: str, filename: str):
    """
    Upload a file to Google Drive (once) and return an embed URL for the
    corresponding Google App viewer.  Subsequent calls for the same unchanged
    file return the cached URL instantly — no re-upload.
    """
    _safe_filename(filename)
    base = FOLDER_MAP.get(folder)
    if not base:
        raise HTTPException(400, "Unknown folder")
    f = (base / filename).resolve()
    if not str(f).startswith(str(base.resolve())) or not f.exists():
        raise HTTPException(404, "Not found")

    suffix = f.suffix.lower()
    if suffix not in _TO_GOOGLE_MIME:
        raise HTTPException(400, f"Unsupported type for Google viewer: {suffix}")

    # ── credentials check ────────────────────────────────────────────────────
    if not Path(_GDRIVE_SECRETS).exists():
        return JSONResponse(status_code=503, content={
            "error": "not_configured",
            "message": "client_secrets.json not found in the Athena root folder.",
        })
    _svc = _build_drive_service()
    if _svc is None:
        return JSONResponse(status_code=503, content={
            "error": "not_configured",
            "needs_auth": True,
            "message": "Google Drive not authorized. Click 'Connect Google Drive' to complete setup.",
        })

    # ── cache lookup (skip upload if file unchanged) ─────────────────────────
    cache_key  = f"{folder}/{filename}"
    file_mtime = str(f.stat().st_mtime)
    cache      = _drive_cache_load()
    entry      = cache.get(cache_key, {})

    if entry.get("mtime") == file_mtime and entry.get("file_id"):
        google_mime = entry.get("google_mime")
        embed_url   = _EMBED_URL.get(google_mime, _EMBED_URL[None]).format(id=entry["file_id"])
        return {"url": embed_url, "file_id": entry["file_id"], "cached": True}

    # ── upload to Drive ───────────────────────────────────────────────────────
    service = _svc

    upload_mime, google_mime = _TO_GOOGLE_MIME[suffix]
    file_metadata = {"name": filename}
    if _GDRIVE_FOLDER_ID:
        file_metadata["parents"] = [_GDRIVE_FOLDER_ID]
    if google_mime:
        file_metadata["mimeType"] = google_mime  # convert to native Docs/Slides

    from googleapiclient.http import MediaFileUpload
    media = MediaFileUpload(str(f), mimetype=upload_mime, resumable=False)

    # Delete the previous Drive copy (avoid accumulating stale versions)
    if entry.get("file_id"):
        try:
            service.files().delete(fileId=entry["file_id"]).execute()
        except Exception:
            pass

    uploaded    = service.files().create(
        body=file_metadata, media_body=media, fields="id,mimeType"
    ).execute()
    file_id     = uploaded["id"]
    actual_mime = uploaded.get("mimeType") or (google_mime or upload_mime)

    cache[cache_key] = {
        "file_id": file_id, "mtime": file_mtime, "google_mime": actual_mime,
    }
    _drive_cache_save(cache)

    embed_url = _EMBED_URL.get(actual_mime, _EMBED_URL[None]).format(id=file_id)
    return {"url": embed_url, "file_id": file_id, "cached": False}


# ── Human Review Queue (Phase 1A) ────────────────────────────────────────────

@app.get("/api/review/pending")
def review_pending():
    if not mem: return {"items": []}
    return {"items": mem.get_pending_reviews(), "stats": mem.get_review_stats()}

async def _resume_workflow(item_id: int, decision: str, notes: str):
    """If the reviewed item belongs to a paused workflow, resume it."""
    if not mem:
        return
    row = mem.get_review(item_id)
    thread_id = row.get("thread_id") if row else None
    if not thread_id:
        return  # standalone item (e.g. legacy subprocess brief) — nothing to resume
    try:
        result = await asyncio.to_thread(_orchestrator().resume_coaching, thread_id, decision, notes)
        for step in result.get("steps", []):
            await manager.broadcast({"type": "agent_log", "agent": "orchestrator", "line": step})
        # Orchestrator returns "finalized" (approved+delivered) or "rejected" — both are
        # successful terminal states. Map to "success" so the UI doesn't show "failed".
        orch_status = result.get("status", "")
        ws_status = "success" if orch_status in ("finalized", "approved", "rejected") else "error"
        await manager.broadcast({"type": "agent_done", "agent": "orchestrator",
                                 "status": ws_status, "ts": datetime.now().isoformat()})
    except Exception as e:
        await manager.broadcast({"type": "agent_log", "agent": "orchestrator",
                                 "line": f"[ERROR] resume failed: {e}"})
        await manager.broadcast({"type": "agent_done", "agent": "orchestrator",
                                 "status": "error", "ts": datetime.now().isoformat()})

@app.post("/api/review/{item_id}/approve")
async def review_approve(item_id: int, request: Request):
    body  = await request.json()
    notes = body.get("notes", "")
    row   = mem.get_review(item_id) if mem else None
    if mem:
        mem.approve_review(item_id, notes)
        _save_review_annotation(item_id, "approved", notes)
    await _resume_workflow(item_id, "approved", notes)
    # Broadcast so the correct module refreshes to show the newly approved deliverable.
    await manager.broadcast({
        "type":      "doc_shipped",
        "id":        item_id,
        "agent":     row.get("agent")     if row else None,
        "item_type": row.get("item_type") if row else None,
        "title":     row.get("title")     if row else None,
        "ts":        datetime.now().isoformat(),
    })
    return {"status": "approved", "id": item_id}

@app.post("/api/review/{item_id}/reject")
async def review_reject(item_id: int, request: Request):
    body  = await request.json()
    notes = body.get("notes", "")
    if mem:
        mem.reject_review(item_id, notes)
        _save_review_annotation(item_id, "rejected", notes)
    await _resume_workflow(item_id, "rejected", notes)
    return {"status": "rejected", "id": item_id}


def _save_review_annotation(item_id: int, decision: str, notes: str):
    """Persist Steven's review decision as a KB annotation for agent grounding."""
    if not mem or not notes.strip():
        return
    try:
        item = mem.get_review(item_id)
        if not item:
            return
        # Infer regulatory domain from agent/item_type
        domain_map = {
            "regulatory_strategy_agent": "Regulatory",
            "sow_agent":                 "Consulting",
            "coaching_brief":            "Coaching",
            "briefing_agent":            "Briefing",
            "deck_agent":                "Consulting",
            "rag_agent":                 "Regulatory",
        }
        domain = domain_map.get(item.get("agent", ""), "General")

        # Write annotation to KB Annotations dir so it's indexed on next KB load
        ann_dir = ATHENA / "knowledge_base" / "Annotations"
        ann_dir.mkdir(parents=True, exist_ok=True)
        import hashlib as _hl
        ann_id   = _hl.md5(f"{item_id}{decision}".encode()).hexdigest()[:10]
        ann_file = ann_dir / f"annotation_{ann_id}.json"
        ann_data = {
            "doc_id":    f"annotation_{ann_id}",
            "doc_name":  f"Steven review: {item.get('title', 'Unknown')}",
            "category":  "Annotations",
            "source":    "steven_review",
            "ingested_at": datetime.now().isoformat(),
            "chunk_count": 1,
            "chunks": [{
                "id":        f"annotation_{ann_id}_chunk_0",
                "doc_id":    f"annotation_{ann_id}",
                "doc_name":  f"Steven review: {item.get('title', '')}",
                "category":  "Annotations",
                "chunk_num": 0,
                "text": (
                    f"Steven Tran review decision: {decision.upper()}\n"
                    f"Item: {item.get('title', '')}\n"
                    f"Agent: {item.get('agent', '')}\n"
                    f"Type: {item.get('item_type', '')}\n"
                    f"Review notes: {notes}\n"
                    f"Domain: {domain}"
                ),
                "word_count": len(notes.split()) + 20,
            }],
        }
        ann_file.write_text(
            __import__('json').dumps(ann_data, indent=2, ensure_ascii=False),
            encoding="utf-8"
        )
        mem.save_annotation(
            review_item_id=item_id,
            decision=decision,
            agent=item.get("agent", ""),
            item_type=item.get("item_type", ""),
            title=item.get("title", ""),
            notes=notes,
            annotation_path=str(ann_file),
            regulatory_domain=domain,
        )
    except Exception as e:
        import logging as _log
        _log.getLogger("server").warning(f"Annotation save failed for review {item_id}: {e}")

@app.get("/api/review/history")
def review_history():
    if not mem: return {"items": []}
    return {"items": mem.get_reviewed_items()}

@app.post("/api/review/{item_id}/reopen")
def review_reopen(item_id: int):
    if mem: mem.reopen_review(item_id)
    return {"status": "pending", "id": item_id}


@app.get("/api/review/{item_id}/summary")
def review_summary(item_id: int):
    """Return a concise 2-3 sentence spoken summary of a review item (voice-ready)."""
    if not mem:
        return JSONResponse(status_code=404, content={"error": "Memory not available"})
    row = mem.get_review(item_id)
    if not row:
        return JSONResponse(status_code=404, content={"error": "Not found"})

    title     = row.get("title", "untitled")
    item_type = row.get("item_type", "")

    # Read file content (markdown / text only)
    raw = ""
    f   = _resolve_review_file(item_id)
    if f and f.suffix.lower() in (".md", ".txt"):
        try:
            raw = f.read_text(encoding="utf-8", errors="replace")[:2500]
        except Exception:
            pass

    if not raw:
        return {"summary": title, "title": title, "id": item_id}

    try:
        import anthropic as _ant
        _key = os.getenv("ANTHROPIC_API_KEY", "")
        if not _key:
            raise RuntimeError("no key")
        _client = _ant.Anthropic(api_key=_key)
        _resp = _client.messages.create(
            model="claude-haiku-4-5",
            max_tokens=120,
            system=(
                "Summarise this document for spoken voice delivery. "
                "Write two or three sentences of plain prose — no markdown, no bullets, no headings. "
                "State the key finding, recommendation, or what was produced. Be specific and concise."
            ),
            messages=[{"role": "user",
                       "content": f"Title: {title}\nType: {item_type}\n\n{raw}"}],
        )
        summary = _resp.content[0].text.strip()
    except Exception:
        summary = title

    return {"summary": summary, "title": title, "id": item_id}


def _resolve_review_file(item_id: int):
    """Locate the file backing a review item by its stored path, falling back
    to a filename search across the known folders if the absolute path moved."""
    if not mem:
        return None
    row = mem.get_review(item_id)
    if not row or not row.get("file_path"):
        return None
    f = Path(row["file_path"])
    if f.exists():
        return f
    for base in FOLDER_MAP.values():
        cand = base / f.name
        if cand.exists():
            return cand
    return None


@app.get("/api/review/{item_id}/content")
def review_content(item_id: int):
    """Return the text of a review item so it can be opened/read in-queue.
    Markdown/text is returned inline; binary docs point the client at /serve."""
    f = _resolve_review_file(item_id)
    if not f:
        return JSONResponse(status_code=404, content={"error": "File not found for this review item"})
    ext = f.suffix.lower().lstrip(".")
    if ext in ("md", "txt"):
        return {"filename": f.name, "ext": ext,
                "content": f.read_text(encoding="utf-8", errors="replace")}
    return {"filename": f.name, "ext": ext, "content": None}   # docx/pdf → use /serve


@app.get("/api/review/{item_id}/serve")
def review_serve(item_id: int):
    """Serve raw bytes of a review item's file (for docx/pdf in-browser viewing)."""
    f = _resolve_review_file(item_id)
    if not f:
        return JSONResponse(status_code=404, content={"error": "File not found"})
    suffix = f.suffix.lower()
    mime = {
        ".pdf":  "application/pdf",
        ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ".md":   "text/plain; charset=utf-8",
        ".txt":  "text/plain; charset=utf-8",
    }.get(suffix, "application/octet-stream")
    return FResponse(content=f.read_bytes(), media_type=mime,
                     headers={"Content-Disposition": f'inline; filename="{f.name}"'})


@app.get("/api/review/{item_id}/google-view")
def review_google_view(item_id: int):
    """
    Upload a review-queue file to Google Drive (once, mtime-cached) and return an
    embed URL for the Google Docs / Slides / Drive viewer.  Falls back gracefully if
    Drive is not yet authorized.
    """
    f = _resolve_review_file(item_id)
    if not f:
        return JSONResponse(status_code=404, content={"error": "File not found"})

    suffix = f.suffix.lower()
    if suffix not in _TO_GOOGLE_MIME:
        return JSONResponse(status_code=400, content={
            "error": "unsupported",
            "message": f"No Google viewer for {suffix} files.",
        })

    # ── credentials check ────────────────────────────────────────────────────
    if not Path(_GDRIVE_SECRETS).exists():
        return JSONResponse(status_code=503, content={
            "error": "not_configured",
            "message": "client_secrets.json not found in the Athena root folder.",
        })
    svc = _build_drive_service()
    if svc is None:
        return JSONResponse(status_code=503, content={
            "error": "not_configured",
            "needs_auth": True,
            "message": "Google Drive not authorized. Click 'Connect Google Drive' to complete setup.",
        })

    # ── cache lookup ─────────────────────────────────────────────────────────
    cache_key  = f"__review__/{item_id}/{f.name}"
    file_mtime = str(f.stat().st_mtime)
    cache      = _drive_cache_load()
    entry      = cache.get(cache_key, {})

    if entry.get("mtime") == file_mtime and entry.get("file_id"):
        google_mime = entry.get("google_mime")
        embed_url   = _EMBED_URL.get(google_mime, _EMBED_URL[None]).format(id=entry["file_id"])
        return {"url": embed_url, "file_id": entry["file_id"], "cached": True}

    # ── upload ───────────────────────────────────────────────────────────────
    upload_mime, google_mime = _TO_GOOGLE_MIME[suffix]
    file_metadata = {"name": f.name}
    if _GDRIVE_FOLDER_ID:
        file_metadata["parents"] = [_GDRIVE_FOLDER_ID]
    if google_mime:
        file_metadata["mimeType"] = google_mime

    from googleapiclient.http import MediaFileUpload
    media = MediaFileUpload(str(f), mimetype=upload_mime, resumable=False)

    if entry.get("file_id"):
        try:
            svc.files().delete(fileId=entry["file_id"]).execute()
        except Exception:
            pass

    uploaded    = svc.files().create(
        body=file_metadata, media_body=media, fields="id,mimeType"
    ).execute()
    file_id     = uploaded["id"]
    actual_mime = uploaded.get("mimeType") or (google_mime or upload_mime)

    cache[cache_key] = {"file_id": file_id, "mtime": file_mtime, "google_mime": actual_mime}
    _drive_cache_save(cache)

    embed_url = _EMBED_URL.get(actual_mime, _EMBED_URL[None]).format(id=file_id)
    return {"url": embed_url, "file_id": file_id, "cached": False}


def _split_frontmatter(text: str):
    """Return (frontmatter_block_including_fences, body). Empty fm if none."""
    if text.startswith("---"):
        end = text.find("\n---", 3)
        if end != -1:
            nl = text.find("\n", end + 1)
            if nl != -1:
                return text[:nl + 1], text[nl + 1:]
    return "", text


def _ai_edit_document(body: str, instruction: str) -> str:
    """Revise a Markdown document per a reviewer instruction, at consulting
    quality, returning the full revised Markdown body (no preamble)."""
    import anthropic
    key = os.getenv("ANTHROPIC_API_KEY", "")
    if not key:
        raise RuntimeError("ANTHROPIC_API_KEY not configured")
    client = anthropic.Anthropic(api_key=key)
    system = (
        "You are a senior editor at a Big 4-calibre MedTech/Pharma consulting firm "
        "(Latitude MedTech). Revise the document per the reviewer's instruction while "
        "preserving its Markdown structure and factual claims. Keep the house style: "
        "specific citations (21 CFR §, ISO clause, company, dollar figure), no filler "
        "phrases ('it is important to note', 'leverage', 'robust', 'in conclusion'), "
        "concrete opening sentence. Return ONLY the full revised Markdown document — "
        "no commentary, no code fences, no preamble."
    )
    resp = client.messages.create(
        model="claude-sonnet-4-6", max_tokens=4000, system=system,
        messages=[{"role": "user",
                   "content": f"Reviewer instruction:\n{instruction}\n\n"
                              f"---\nDocument to revise:\n\n{body}"}],
    )
    return resp.content[0].text.strip()


class ReviewEditRequest(BaseModel):
    instruction: str


@app.post("/api/review/{item_id}/edit")
async def review_edit(item_id: int, req: ReviewEditRequest):
    """Apply an AI edit to a review item's document from a reviewer prompt.
    Writes the revised content back to the same file and returns it."""
    if not req.instruction.strip():
        return JSONResponse(status_code=400, content={"error": "Instruction is empty"})
    f = _resolve_review_file(item_id)
    if not f:
        return JSONResponse(status_code=404, content={"error": "File not found for this item"})
    if f.suffix.lower() not in (".md", ".txt"):
        return JSONResponse(status_code=400,
            content={"error": f"Only markdown/text documents are editable (got {f.suffix})"})
    try:
        original    = f.read_text(encoding="utf-8", errors="replace")
        fm, body    = _split_frontmatter(original)
        revised     = await asyncio.to_thread(_ai_edit_document, body, req.instruction)
        new_content = (fm + "\n" + revised) if fm else revised
        f.write_text(new_content, encoding="utf-8")
        if mem:
            mem.log_event("human_review", "ai_edit",
                          subject=f.name, metadata={"item_id": item_id})
        return {"status": "edited", "filename": f.name,
                "ext": f.suffix.lower().lstrip("."), "content": new_content}
    except Exception as e:
        return JSONResponse(status_code=500, content={"error": str(e)})


@app.post("/api/documents/edit")
async def edit_document_voice(request: Request, background_tasks: BackgroundTasks):
    """
    AI-edit any document by folder + filename + natural-language instruction.
    Intended for voice agent commands ("add a section on X to the Y report").

    Body: { "folder": str, "filename": str, "instruction": str }

    MD  → read → Claude Sonnet edit → write back in-place.
    DOCX→ extract paragraph text → Claude Sonnet edit → regenerate branded DOCX
          in-place via the existing generate_docx pipeline.

    WebSocket events:
      { type: "doc_edit_progress", filename, folder, message }
      { type: "doc_edit_done",     filename, folder, status: "done"|"error", message? }
    """
    body        = await request.json()
    folder      = body.get("folder", "documents")
    filename    = body.get("filename", "")
    instruction = body.get("instruction", "")

    if not filename or not instruction.strip():
        raise HTTPException(400, "filename and instruction are required")

    _safe_filename(filename)
    base = FOLDER_MAP.get(folder)
    if not base:
        raise HTTPException(400, f"Unknown folder: {folder}")
    f = (base / filename).resolve()
    if not str(f).startswith(str(base.resolve())) or not f.exists():
        raise HTTPException(404, "Document not found")

    async def _run():
        await manager.broadcast({
            "type": "doc_edit_progress",
            "filename": filename, "folder": folder,
            "message": f"Editing \"{filename}\"…",
        })
        try:
            suffix = f.suffix.lower()

            if suffix == ".md":
                original      = f.read_text(encoding="utf-8", errors="replace")
                fm, body_text = _split_frontmatter(original)
                revised       = await asyncio.to_thread(_ai_edit_document, body_text, instruction)
                f.write_text((fm + "\n" + revised) if fm else revised, encoding="utf-8")

            elif suffix == ".docx":
                def _extract(path):
                    from docx import Document as _Doc
                    doc, parts = _Doc(str(path)), []
                    for p in doc.paragraphs:
                        if not p.text.strip():
                            continue
                        s = p.style.name if p.style else ""
                        if   "Heading 1" in s: parts.append(f"# {p.text}")
                        elif "Heading 2" in s: parts.append(f"## {p.text}")
                        elif "Heading 3" in s: parts.append(f"### {p.text}")
                        else:                  parts.append(p.text)
                    return "\n\n".join(parts)

                doc_md     = await asyncio.to_thread(_extract, f)
                revised_md = await asyncio.to_thread(_ai_edit_document, doc_md, instruction)
                title      = re.sub(r"^\d{4}[-_]\d{2}[-_]\d{2}[-_]?", "", f.stem).replace("_", " ").title()
                new_path   = await asyncio.to_thread(generate_docx, title, revised_md, "article")
                f.write_bytes(new_path.read_bytes())
                new_path.unlink(missing_ok=True)

            else:
                raise ValueError(f"Unsupported type for editing: {suffix}")

            await manager.broadcast({
                "type": "doc_edit_done",
                "filename": filename, "folder": folder, "status": "done",
            })
        except Exception as exc:
            await manager.broadcast({
                "type": "doc_edit_done",
                "filename": filename, "folder": folder,
                "status": "error", "message": str(exc),
            })

    background_tasks.add_task(_run)
    return {"status": "started", "filename": filename, "folder": folder}


@app.get("/api/dashboard/knowledge-growth")
def knowledge_growth(days: int = 90):
    """Cumulative knowledge accumulation over time — every ingested KB document
    and every learning item, counted per day and accumulated. Powers the
    Workforce 'Knowledge Growth' chart."""
    if not mem:
        return {"daily": [], "total": 0}
    rows = mem.conn.execute(
        """SELECT substr(ts,1,10) AS date, COUNT(*) AS n FROM (
               SELECT timestamp AS ts FROM knowledge_items
               UNION ALL
               SELECT timestamp AS ts FROM agent_learning
           ) WHERE ts IS NOT NULL
           GROUP BY date ORDER BY date ASC"""
    ).fetchall()
    cutoff = (datetime.now() - timedelta(days=days)).strftime("%Y-%m-%d")
    cum, out = 0, []
    for r in rows:
        cum += r["n"]
        if r["date"] >= cutoff:
            out.append({"date": r["date"], "items": r["n"], "cumulative": cum})
    # If the window opens mid-history, seed the first point with the running
    # total so the curve starts at the true accumulated height, not zero.
    return {"daily": out, "total": cum}


# ── Token history for dashboard charts ───────────────────────────────────────

@app.get("/api/dashboard/history")
def dashboard_history(days: int = 30):
    if not mem:
        return {"daily": []}
    cutoff = (datetime.now() - timedelta(days=days)).strftime("%Y-%m-%d")
    rows = mem.conn.execute(
        "SELECT date, total_calls, total_tokens, total_cost, cache_hits "
        "FROM daily_summaries WHERE date >= ? ORDER BY date ASC", (cutoff,)
    ).fetchall()
    return {"daily": [dict(r) for r in rows]}


@app.get("/api/dashboard/timeseries")
def dashboard_timeseries(day: str = "today"):
    """Time-resolved token usage: today vs yesterday totals + per-hour breakdown.

    `day` selects which day the hourly breakdown covers: "today" or "yesterday".
    Timestamps in api_calls are stored as local-time ISO strings (datetime.now().isoformat()),
    so all date math uses SQLite's 'localtime' modifier to stay in the user's timezone.
    """
    if not mem:
        return {"today": {}, "yesterday": {}, "hourly": [], "hourly_day": day}

    def _day_totals(day_expr):
        row = mem.conn.execute(
            "SELECT COUNT(*) AS calls, "
            "       COALESCE(SUM(total_tokens),0) AS tokens, "
            "       COALESCE(SUM(cost_usd),0)     AS cost, "
            "       COALESCE(SUM(cache_hit),0)    AS cache_hits "
            "FROM api_calls WHERE date(timestamp)=" + day_expr
        ).fetchone()
        return dict(row) if row else {}

    today     = _day_totals("date('now','localtime')")
    yesterday = _day_totals("date('now','localtime','-1 day')")
    today["date"]     = datetime.now().strftime("%Y-%m-%d")
    yesterday["date"] = (datetime.now() - timedelta(days=1)).strftime("%Y-%m-%d")

    # Per-hour breakdown for the selected day, zero-filled across all 24 hours.
    hour_day_expr = "date('now','localtime','-1 day')" if day == "yesterday" else "date('now','localtime')"
    rows = mem.conn.execute(
        "SELECT strftime('%H', timestamp) AS hour, "
        "       COUNT(*) AS calls, "
        "       COALESCE(SUM(total_tokens),0) AS tokens, "
        "       COALESCE(SUM(cost_usd),0)     AS cost "
        "FROM api_calls WHERE date(timestamp)=" + hour_day_expr + " "
        "GROUP BY hour"
    ).fetchall()
    by_hour = {r["hour"]: dict(r) for r in rows}
    hourly = [
        by_hour.get(f"{h:02d}", {"hour": f"{h:02d}", "calls": 0, "tokens": 0, "cost": 0})
        for h in range(24)
    ]
    for slot in hourly:
        slot["hour"] = f"{int(slot['hour']):02d}"

    return {"today": today, "yesterday": yesterday, "hourly": hourly, "hourly_day": day}


# ── Voice greeting + shutdown ─────────────────────────────────────────────────

import random as _random

def _time_greeting() -> str:
    h = datetime.now().hour
    if h < 12:   return "Good morning"
    if h < 17:   return "Good afternoon"
    if h < 21:   return "Good evening"
    return "Good evening"

# Note: "Athena" mispronounced as "Atheney" by Kokoro bf_emma.
# Greetings now use "your AI assistant" or "online" — avoids the name in TTS.
# The name still appears visually in the UI.
_GREETINGS_MORNING = [
    "Good morning, Steven. Your AI operating system is online and ready.",
    "Morning. Latitude MedTech is online. What are we working on today?",
    "Good morning, Steven. Your AI assistant has loaded. Ready to go.",
    "Morning. All systems online. What can I help you with?",
]
_GREETINGS_AFTERNOON = [
    "Good afternoon, Steven. Your AI assistant is online and ready.",
    "Afternoon. Latitude MedTech online. What would you like to work on?",
    "Good afternoon, Steven. All systems ready. How can I help?",
    "Afternoon. Online and standing by, Steven.",
]
_GREETINGS_EVENING = [
    "Good evening, Steven. Your AI assistant is online — still at it.",
    "Evening. Latitude MedTech online. What's on the agenda?",
    "Good evening, Steven. Online and ready. What do you need?",
    "Evening. All systems running. Let's get to work.",
]

_GOODBYES = [
    "Goodbye, Steven. Have a wonderful day — you're building something great.",
    "Signing off. Take care, and well done today.",
    "Farewell, Steven. Great working with you. Until next time.",
    "Goodbye. Rest well — Latitude MedTech will be here when you return.",
    "That's a wrap. Have a great day, Steven.",
    "Signing off. It was a pleasure. See you next time.",
]

def _pick_greeting() -> str:
    h = datetime.now().hour
    if h < 12:   pool = _GREETINGS_MORNING
    elif h < 17: pool = _GREETINGS_AFTERNOON
    else:        pool = _GREETINGS_EVENING
    return _random.choice(pool)

def _speak_phrase(text: str):
    """Speak a goodbye phrase via Kokoro TTS with Windows SAPI fallback.

    Synchronous — blocks until TTS completes (needed for the shutdown flow).
    For the startup greeting use _speak_phrase_greeting (non-blocking, longer wait).
    Retries Kokoro once after 4 s then falls back to Windows SAPI.
    """
    import io, json, time, urllib.request as _urllib

    def _try_kokoro() -> bool:
        try:
            import soundfile as sf, sounddevice as sd, numpy as np
            payload = json.dumps(
                {"text": text, "voice": os.getenv("VOICE_KOKORO_VOICE", "bf_emma")}
            ).encode()
            req = _urllib.Request(
                "http://127.0.0.1:8002/speak", data=payload,
                headers={"Content-Type": "application/json"}, method="POST")
            with _urllib.urlopen(req, timeout=20) as resp:
                wav = resp.read()
            data, sr = sf.read(io.BytesIO(wav))
            sd.play(data.astype(np.float32), sr)
            sd.wait()
            return True
        except Exception as exc:
            print(f"[voice] _speak_phrase kokoro failed: {exc!r}", flush=True)
            return False

    if _try_kokoro():
        time.sleep(0.5)
        return

    print("[voice] _speak_phrase: retrying Kokoro in 4 s…", flush=True)
    time.sleep(4)
    if _try_kokoro():
        time.sleep(0.5)
        return

    print("[voice] _speak_phrase: Kokoro unavailable, using Windows SAPI", flush=True)
    try:
        safe  = text.replace("'", "").replace('"', " ")
        script = (
            "Add-Type -AssemblyName System.Speech;"
            "$s=New-Object System.Speech.Synthesis.SpeechSynthesizer;"
            f"$s.Rate=0;$s.Speak('{safe}')"
        )
        flags = getattr(subprocess, "CREATE_NO_WINDOW", 0)
        subprocess.run(
            ["powershell.exe", "-NoProfile", "-NonInteractive", "-Command", script],
            timeout=30, creationflags=flags, stdin=subprocess.DEVNULL,
        )
    except Exception as exc:
        print(f"[voice] _speak_phrase SAPI failed: {exc!r}", flush=True)
    time.sleep(0.5)


def _speak_phrase_greeting(text: str):
    """Greeting TTS: waits up to 150 s for Kokoro to finish its model load.

    Cold-start Kokoro takes 60–90 s to load; this function polls every 3 s so
    the greeting plays via the high-quality Kokoro voice once it is ready.
    Falls back to Windows SAPI if Kokoro never comes up within the window.
    Runs in a background executor so the greet endpoint returns immediately.
    """
    import io, json, time, urllib.request as _urllib

    deadline = time.monotonic() + 150
    while time.monotonic() < deadline:
        try:
            import soundfile as sf, sounddevice as sd, numpy as np
            payload = json.dumps(
                {"text": text, "voice": os.getenv("VOICE_KOKORO_VOICE", "bf_emma")}
            ).encode()
            req = _urllib.Request(
                "http://127.0.0.1:8002/speak", data=payload,
                headers={"Content-Type": "application/json"}, method="POST")
            with _urllib.urlopen(req, timeout=20) as resp:
                wav = resp.read()
            data, sr = sf.read(io.BytesIO(wav))
            sd.play(data.astype(np.float32), sr)
            sd.wait()
            print("[voice] greeting: played via Kokoro", flush=True)
            return
        except Exception as exc:
            remaining = deadline - time.monotonic()
            print(f"[voice] greeting: Kokoro not ready ({exc!r}), {remaining:.0f}s remaining…", flush=True)
            if remaining > 3:
                time.sleep(3)

    # Kokoro never came up — SAPI fallback
    print("[voice] greeting: Kokoro timed out after 150 s, using Windows SAPI", flush=True)
    try:
        safe = text.replace("'", "").replace('"', " ")
        script = (
            "Add-Type -AssemblyName System.Speech;"
            "$s=New-Object System.Speech.Synthesis.SpeechSynthesizer;"
            f"$s.Rate=0;$s.Speak('{safe}')"
        )
        flags = getattr(subprocess, "CREATE_NO_WINDOW", 0)
        subprocess.run(
            ["powershell.exe", "-NoProfile", "-NonInteractive", "-Command", script],
            timeout=30, creationflags=flags, stdin=subprocess.DEVNULL,
        )
    except Exception as exc:
        print(f"[voice] greeting SAPI failed: {exc!r}", flush=True)


_session_greeted = False   # plays once per server restart, regardless of tab refreshes

@app.post("/api/voice/greet")
async def greet():
    """
    Play the startup greeting exactly once per server session.
    React StrictMode and tab refreshes call this multiple times — the flag prevents repeats.

    Fires TTS in a background executor so the HTTP response returns immediately —
    Kokoro cold-start takes 60–90 s and we must not hold the frontend.
    """
    global _session_greeted
    if _session_greeted:
        return {"phrase": "", "skipped": True}
    _session_greeted = True
    phrase = _pick_greeting()
    # Non-blocking: greeting waits for Kokoro in the background.
    asyncio.ensure_future(asyncio.get_event_loop().run_in_executor(
        None, _speak_phrase_greeting, phrase
    ))
    return {"phrase": phrase}

@app.post("/api/shutdown")
async def shutdown_app(background_tasks: BackgroundTasks):
    """Play a goodbye, hand off to the verified stop script, then exit.

    TTS is awaited synchronously so the response only returns after Athena
    finishes speaking. The stop script and os._exit run as a background task
    so they fire after the response has been fully sent to the client.
    """
    global _shutting_down
    if _shutting_down:             # guard: double-click of the off button sends two POSTs
        return {"status": "already_shutting_down"}
    _shutting_down = True          # prevent WebSocket-disconnect goodbye from double-firing
    phrase = _random.choice(_GOODBYES)
    stop_script = ATHENA / "ui" / "stop_athena.ps1"

    # Block until TTS finishes — response is sent only after this completes.
    await asyncio.get_event_loop().run_in_executor(None, _speak_phrase, phrase)

    def _cleanup():
        import time; time.sleep(0.15)  # brief buffer for response to flush
        try:
            flags = getattr(subprocess, "CREATE_NO_WINDOW", 0) | getattr(subprocess, "CREATE_NEW_PROCESS_GROUP", 0)
            subprocess.Popen(
                ["powershell.exe", "-ExecutionPolicy", "Bypass", "-NoProfile",
                 "-File", str(stop_script)],
                creationflags=flags, close_fds=True,
            )
        except Exception:
            pass
        import os; os._exit(0)

    background_tasks.add_task(_cleanup)
    return {"phrase": phrase, "status": "shutting_down"}


# ── Security status endpoint (SOC II CC7.2 — system monitoring) ───────────────

@app.get("/api/security/status")
def security_status():
    """Returns a summary of active security controls for monitoring purposes."""
    import stat as _stat
    def _file_perm(p: Path) -> str:
        try:
            mode = p.stat().st_mode
            owner_r  = bool(mode & _stat.S_IRUSR)
            owner_w  = bool(mode & _stat.S_IWUSR)
            group_r  = bool(mode & _stat.S_IRGRP)
            other_r  = bool(mode & _stat.S_IROTH)
            return f"rw={owner_r and owner_w}, group_read={group_r}, other_read={other_r}"
        except Exception:
            return "unknown"

    db_path  = ATHENA / "memory" / "latitude_memory.db"
    key_path = _KEY_FILE
    audit_lines = 0
    try:
        with open(_AUDIT_LOG, "r", encoding="utf-8") as f:
            audit_lines = sum(1 for _ in f)
    except Exception:
        pass

    return {
        "controls": {
            "security_headers":      True,
            "cors_whitelist":        True,
            "rate_limiting":         "120 req/min",
            "request_body_limit":    f"{_MAX_BODY_BYTES // (1024*1024)} MB",
            "path_traversal_guard":  True,
            "shell_injection_guard": True,
            "session_token":         True,
            "audit_logging":         True,
        },
        "file_permissions": {
            "database": _file_perm(db_path)  if db_path.exists()  else "missing",
            "key_file":  _file_perm(key_path) if key_path.exists() else "missing",
        },
        "audit_log": {
            "path":       str(_AUDIT_LOG),
            "total_entries": audit_lines,
        },
        "gaps": [
            "HTTPS not enabled (localhost deployment — Electron shell handles TLS boundary)",
            "API key auth not enforced on all routes (Electron shell is trusted origin via CORS)",
            "Database at rest not encrypted (SQLite; file permissions restrict access to current OS user)",
        ],
        "standard": "SOC II Type II relevant controls — CC6.1, CC6.6, CC6.7, CC7.2, A1.2",
    }


# ── Phase 2C: Client Lifecycle API ───────────────────────────────────────────

class ClientCreateRequest(BaseModel):
    name:                 str
    org:                  str  = ""
    role:                 str  = ""
    email:                str  = ""
    phone:                str  = ""
    program_tier:         str  = ""
    regulatory_challenge: str  = ""
    timeline:             str  = ""
    budget_range:         str  = ""
    status:               str  = "prospect"
    source_target_id:     int  = None
    notes:                str  = ""

class ClientUpdateRequest(BaseModel):
    org:                  str  = None
    role:                 str  = None
    email:                str  = None
    phone:                str  = None
    program_tier:         str  = None
    regulatory_challenge: str  = None
    timeline:             str  = None
    budget_range:         str  = None
    status:               str  = None
    notes:                str  = None

class EngagementCreateRequest(BaseModel):
    title:       str
    description: str  = ""
    status:      str  = "scoping"
    start_date:  str  = None
    end_date:    str  = None
    value_usd:   float= None
    notes:       str  = ""

class EngagementUpdateRequest(BaseModel):
    title:       str   = None
    description: str   = None
    sow_path:    str   = None
    status:      str   = None
    start_date:  str   = None
    end_date:    str   = None
    value_usd:   float = None
    notes:       str   = None

class SOWTriggerRequest(BaseModel):
    client_id:    int
    engagement_id:int = None

class RegulatoryAssessmentRequest(BaseModel):
    device_type:    str
    classification: str
    markets:        list = ["US"]
    qms_state:      str
    context:        str  = ""
    client_id:      int  = None

@app.get("/api/clients")
def get_clients(status: str = None):
    if not mem: return {"clients": []}
    return {"clients": mem.get_clients(status=status or None)}

@app.post("/api/clients")
def create_client(req: ClientCreateRequest):
    if not mem: return JSONResponse(status_code=503, content={"error": "Database unavailable — please restart the server."})
    try:
        client_id = mem.add_client(
            name=req.name, org=req.org or None, role=req.role or None,
            email=req.email or None, phone=req.phone or None,
            program_tier=req.program_tier or None,
            regulatory_challenge=req.regulatory_challenge or None,
            timeline=req.timeline or None, budget_range=req.budget_range or None,
            status=req.status, source_target_id=req.source_target_id,
            notes=req.notes or None,
        )
        if not client_id:
            return JSONResponse(status_code=500, content={"error": "Insert returned no ID — check the database."})
        return {"status": "created", "client_id": client_id}
    except Exception as exc:
        return JSONResponse(status_code=500, content={"error": str(exc)})

@app.get("/api/clients/{client_id}")
def get_client(client_id: int):
    if not mem: return JSONResponse(status_code=503, content={"error": "DB unavailable"})
    client = mem.get_client(client_id)
    if not client: return JSONResponse(status_code=404, content={"error": "Not found"})
    return client

@app.put("/api/clients/{client_id}")
async def update_client(client_id: int, req: ClientUpdateRequest):
    if not mem: return JSONResponse(status_code=503, content={"error": "DB unavailable"})
    fields = {k: v for k, v in req.model_dump().items() if v is not None}
    mem.update_client(client_id, **fields)
    return {"status": "updated", "client_id": client_id}

@app.delete("/api/clients/{client_id}")
def delete_client(client_id: int):
    if not mem: return JSONResponse(status_code=503, content={"error": "DB unavailable"})
    mem.delete_client(client_id)
    return {"status": "deleted", "client_id": client_id}

@app.get("/api/clients/{client_id}/engagements")
def get_engagements(client_id: int):
    if not mem: return {"engagements": []}
    return {"engagements": mem.get_engagements(client_id)}

@app.post("/api/clients/{client_id}/engagements")
def create_engagement(client_id: int, req: EngagementCreateRequest):
    if not mem: return JSONResponse(status_code=503, content={"error": "DB unavailable"})
    eng_id = mem.add_engagement(
        client_id=client_id, title=req.title,
        description=req.description or None, status=req.status,
        start_date=req.start_date, end_date=req.end_date,
        value_usd=req.value_usd, notes=req.notes or None,
    )
    return {"status": "created", "engagement_id": eng_id}

@app.put("/api/engagements/{engagement_id}")
async def update_engagement(engagement_id: int, req: EngagementUpdateRequest):
    if not mem: return JSONResponse(status_code=503, content={"error": "DB unavailable"})
    fields = {k: v for k, v in req.model_dump().items() if v is not None}
    mem.update_engagement(engagement_id, **fields)
    return {"status": "updated", "engagement_id": engagement_id}

@app.post("/api/agents/sow")
async def trigger_sow(req: SOWTriggerRequest, background_tasks: BackgroundTasks):
    args = ["--client-id", str(req.client_id)]
    if req.engagement_id:
        args += ["--engagement-id", str(req.engagement_id)]
    background_tasks.add_task(run_agent, "sow_agent", "sow_agent.py", args)
    return {"status": "started", "agent": "sow_agent", "client_id": req.client_id}

@app.post("/api/agents/regulatory-assessment")
async def trigger_regulatory_assessment(req: RegulatoryAssessmentRequest,
                                         background_tasks: BackgroundTasks):
    args = [
        "--device-type",    _safe_arg(req.device_type),
        "--classification", _safe_arg(req.classification),
        "--markets",        *[_safe_arg(m) for m in req.markets],
        "--qms-state",      _safe_arg(req.qms_state),
    ]
    if req.context:
        args += ["--context", _safe_arg(req.context)]
    if req.client_id:
        args += ["--client-id", str(req.client_id)]
    background_tasks.add_task(
        run_agent, "regulatory_strategy_agent", "regulatory_strategy_agent.py", args,
        req.device_type
    )
    return {"status": "started", "agent": "regulatory_strategy_agent",
            "device_type": req.device_type}

@app.get("/api/annotations")
def get_annotations(decision: str = None, limit: int = 50):
    if not mem: return {"annotations": []}
    return {"annotations": mem.get_annotations(decision=decision, limit=limit)}


if __name__ == "__main__":
    import uvicorn
    print("\nLatitude MedTech API starting on http://localhost:8000")
    print("Open http://localhost:3000 for the dashboard\n")
    uvicorn.run(app, host="127.0.0.1", port=8000, log_level="warning")
